#!/usr/bin/env python3
"""Builds the ARE Mid-Term presentation deck (16:9, >=10 slides)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
FIGS = os.path.join(HERE, "figs")
OUT = os.path.join(os.path.dirname(HERE), "ARE_MidTerm_Presentation.pptx")

PRIMARY = RGBColor(0x1F, 0x4E, 0x79)
ACCENT = RGBColor(0x2E, 0x75, 0xB6)
DARK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEA, 0xF1, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def _set(run, size, bold=False, color=DARK, italic=False, font="Calibri"):
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font

def rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf

def footer(slide, n):
    tb, tf = textbox(slide, Inches(0.4), Inches(7.05), Inches(9), Inches(0.35))
    r = tf.paragraphs[0].add_run()
    r.text = "Adaptive Rendering Engine  |  Mid-Term Presentation"
    _set(r, 10, color=GRAY)
    tb2, tf2 = textbox(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35))
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r2 = p.add_run(); r2.text = str(n); _set(r2, 10, color=GRAY)

def content_slide(title, n, band_h=1.05):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, Inches(band_h), PRIMARY)
    rect(s, 0, Inches(band_h), SW, Inches(0.06), ACCENT)
    tb, tf = textbox(s, Inches(0.5), Inches(0.18), Inches(12.3), Inches(band_h - 0.3))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    r = tf.paragraphs[0].add_run(); r.text = title; _set(r, 30, True, WHITE)
    footer(s, n)
    return s

def bullets(slide, items, left=0.7, top=1.5, width=12.0, height=5.2, size=20, gap=10):
    tb, tf = textbox(slide, Inches(left), Inches(top), Inches(width), Inches(height))
    for i, item in enumerate(items):
        lvl = 0
        text = item
        if isinstance(item, tuple):
            lvl, text = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        bullet = "•  " if lvl == 0 else "–  "
        r = p.add_run(); r.text = ("    " * lvl) + bullet + text
        _set(r, size if lvl == 0 else size - 2, bold=(lvl == 0 and False), color=DARK)
    return tb

def image_slide(title, n, img, caption=None, img_w=8.6, top=1.45):
    s = content_slide(title, n)
    path = os.path.join(FIGS, img)
    pic = s.shapes.add_picture(path, Inches(0), Inches(top), width=Inches(img_w))
    pic.left = int((SW - pic.width) / 2)
    if caption:
        cy = Inches(top) + pic.height + Inches(0.1)
        tb, tf = textbox(s, Inches(1), cy, Inches(11.3), Inches(0.5))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = caption; _set(r, 14, italic=True, color=GRAY)
    return s

def table_slide(title, n, headers, rows, col_w, top=1.5, fsize=15, hfsize=15):
    s = content_slide(title, n)
    nrows, ncols = len(rows) + 1, len(headers)
    total_w = sum(col_w)
    left = int((SW - Inches(total_w)) / 2)
    gtab = s.shapes.add_table(nrows, ncols, left, Inches(top),
                              Inches(total_w), Inches(0.4 * nrows)).table
    for j, w in enumerate(col_w):
        gtab.columns[j].width = Inches(w)
    for j, h in enumerate(headers):
        c = gtab.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = PRIMARY
        para = c.text_frame.paragraphs[0]
        _set(para.runs[0], hfsize, True, WHITE)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gtab.cell(i, j); c.text = str(val)
            c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE
            _set(c.text_frame.paragraphs[0].runs[0], fsize, color=DARK)
    return s

# ---------------------------------------------------------------- slide 1 ----
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, 0, Inches(2.55), SW, Inches(0.06), ACCENT)
tb, tf = textbox(s, Inches(0.8), Inches(1.0), Inches(11.7), Inches(1.6))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "ADAPTIVE RENDERING ENGINE"; _set(r, 44, True, WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "A Runtime that Selects the Optimal Web Rendering Strategy per Request"
_set(r, 18, color=LIGHT, italic=True)
tb, tf = textbox(s, Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.4))
def cl(tf, text, size, bold=False, color=WHITE, first=False, after=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(after)
    r = p.add_run(); r.text = text; _set(r, size, bold, color)
cl(tf, "Mid-Term Project Presentation", 20, True, LIGHT, first=True, after=14)
cl(tf, "Submitted by", 15, color=LIGHT)
cl(tf, "Bijay Bk (220305)   •   Devendra Pandey (220306)", 17, True)
cl(tf, "Manish Joshi (220312)   •   Pramod Panta (220317)", 17, True, after=14)
cl(tf, "Under the Supervision of", 15, color=LIGHT)
cl(tf, "Er. Robinhood Khadka", 18, True, after=14)
cl(tf, "Department of ICT and Computer Engineering", 14, color=LIGHT)
cl(tf, "Cosmos College of Management & Technology (Affiliated to Pokhara University)", 14, color=LIGHT)
cl(tf, "Bachelor of Engineering in Computer  •  June 2026", 14, color=LIGHT)

# ---------------------------------------------------------------- slide 2 ----
s = content_slide("Introduction", 2)
bullets(s, [
    "The Adaptive Rendering Engine (ARE) is a runtime that automatically "
    "selects the best web rendering strategy for every request.",
    "It chooses among six strategies: SSG, SSR, Streaming SSR, ISR, CSR and "
    "simulated Edge-ISR.",
    "The decision is based on live context: network speed, device type, cache "
    "freshness, server load and data volatility.",
    "Built as technology (an engine) — not an application — using only "
    "open-source tools at zero budget.",
    "Runs locally on a Docker-based private server; no commercial cloud needed.",
], size=21, gap=14)

# ---------------------------------------------------------------- slide 3 ----
s = content_slide("Problem & Motivation", 3)
bullets(s, [
    "Frameworks like Next.js and Remix support many rendering strategies …",
    (1, "but the strategy is fixed by the developer at build time."),
    (1, "the same strategy is applied to every user, regardless of context."),
    "Poor, static choices lead to:",
    (1, "high Time-to-First-Byte (TTFB) and excessive JavaScript hydration;"),
    (1, "increased server load and poor experience on low-end devices."),
    "Core problem: there is no automated system that selects the optimal "
    "rendering strategy at runtime based on real-time context.",
], size=20, gap=10)

# ---------------------------------------------------------------- slide 4 ----
s = content_slide("Objectives", 4)
bullets(s, [
    "General Objective:",
    (1, "Design, implement and evaluate an engine that selects rendering "
        "strategies at runtime to optimise web performance."),
    "Specific Objectives:",
    (1, "Define the contextual variables (network, load, device, cache, volatility)."),
    (1, "Build a pure, rule-based decision engine mapping context → strategy."),
    (1, "Implement six modular strategies behind one interface."),
    (1, "Measure TTFB, render time, response size, cache efficiency, resources."),
    (1, "Compare adaptive vs. static rendering on a Docker private server."),
], size=19, gap=8)

# ---------------------------------------------------------------- slide 5 ----
image_slide("System Architecture", 5, "system-architecture.png",
            "Origin + two edge nodes + reverse proxy + shared cache on one private Docker network",
            img_w=8.2, top=1.5)

# ---------------------------------------------------------------- slide 6 ----
s = content_slide("Decision Engine", 6)
pic = s.shapes.add_picture(os.path.join(FIGS, "decision-flow.png"), 0, Inches(1.4), height=Inches(5.4))
pic.left = Inches(0.6)
bullets(s, [
    "Pure, deterministic function.",
    "No I/O — fully unit-testable.",
    "Rules evaluated top-to-bottom.",
    "First match wins.",
    "Unconditional fallback → SSR.",
    "All 9 rules covered by tests.",
], left=7.4, top=1.7, width=5.4, size=19, gap=14)

# ---------------------------------------------------------------- slide 7 ----
s = content_slide("Rendering Pipeline & Strategies", 7)
pic = s.shapes.add_picture(os.path.join(FIGS, "rendering-pipeline.png"), 0, Inches(1.55), width=Inches(10.4))
pic.left = int((SW - pic.width) / 2)
bullets(s, [
    "SSG – pre-built HTML  •  SSR – fresh per request  •  Streaming SSR – shell first, stream the rest",
    "ISR – cached + revalidate  •  CSR – client renders  •  Edge-ISR – cached at the edge",
], left=0.7, top=5.7, width=12.0, size=16, gap=8)

# ---------------------------------------------------------------- slide 8 ----
s = content_slide("Technology Stack & Private Server", 8)
bullets(s, [
    "Language & runtime: Node.js 20+ with TypeScript.",
    "View layer: React 18 (SSR, streaming, hydration) — used as the rendering primitive.",
    "Server: native Node HTTP (the engine is the runtime; no framework lock-in).",
    "Tooling: esbuild (client bundle), Vitest (tests), Apache Bench (load).",
    "Caching: filesystem + memory + optional Redis (shared edge cache).",
    "Private server: Docker Compose — origin + 2 edge nodes + nginx proxy + Redis.",
], size=20, gap=12)

# ---------------------------------------------------------------- slide 9 ----
s = content_slide("Progress So Far", 9)
bullets(s, [
    "Core runtime complete: context analyzer, decision engine, strategy registry, orchestrator.",
    "All six rendering strategy modules implemented behind one interface.",
    "Three-tier caching with stale-while-revalidate; metrics + report generator.",
    "React test pages + esbuild client bundle (hydration / CSR).",
    "Multi-stage Docker image builds and runs; compose topology defined.",
    "Validation scripts for strategy switching, header inspection and load testing.",
    "Status: compiles cleanly, 14/14 unit tests pass, runs and verified locally.",
], size=19, gap=9)

# ---------------------------------------------------------------- slide 10 ---
table_slide("Preliminary Results — Strategy Switching", 10,
    ["Request context (headers)", "Selected", "Expected"],
    [["static volatility, /static", "SSG", "SSG"],
     ["realtime + fast + desktop, /dynamic", "CSR", "CSR"],
     ["high load, /dynamic", "ISR", "ISR"],
     ["realtime + mobile, /dynamic", "SSR", "SSR"],
     ["heavy payload + medium net, /heavy", "STREAMING_SSR", "STREAMING_SSR"],
     ["static + edge + cold cache, /dynamic", "EDGE_ISR", "EDGE_ISR"],
     ["periodic, /dynamic", "ISR", "ISR"]],
    col_w=[6.2, 2.6, 2.6], top=1.5, fsize=14, hfsize=14)
tf = prs.slides[-1].shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12), Inches(0.5)).text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Same URLs, only headers changed → six different strategies, each matching the expected rule."
_set(r, 15, italic=True, color=ACCENT)

# ---------------------------------------------------------------- slide 11 ---
table_slide("Preliminary Results — Performance", 11,
    ["Strategy", "Avg TTFB (ms)", "Avg Render (ms)", "Cache-hit", "Avg Bytes"],
    [["SSG", "1.24", "0.74", "1.0", "1,199"],
     ["CSR", "0.31", "0.16", "0.0", "401"],
     ["ISR", "0.41", "0.26", "1.0", "1,219"],
     ["SSR", "6.45", "0.16", "0.0", "1,219"],
     ["Streaming SSR", "8.36", "0.79", "0.0", "20,896"],
     ["Edge-ISR", "1.48", "0.20", "1.0", "1,234"]],
    col_w=[2.6, 2.2, 2.4, 1.6, 1.8], top=1.5, fsize=15, hfsize=14)
tf = prs.slides[-1].shapes.add_textbox(Inches(0.7), Inches(6.1), Inches(12), Inches(0.8)).text_frame
tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = ("Trade-offs confirmed: CSR ships the smallest shell (~0.4 KB), Streaming SSR the "
          "largest heavy page (~20 KB); cached strategies (SSG/ISR/Edge-ISR) hit 1.0 with very low TTFB.")
_set(r, 14, italic=True, color=ACCENT)

# ---------------------------------------------------------------- slide 12 ---
table_slide("Challenges & Solutions", 12,
    ["Challenge", "Solution Adopted"],
    [["Non-ASCII in HTTP headers caused errors", "Sanitise header to ASCII; keep full text in logs"],
     ["nginx cannot inject edge latency natively", "Model edges as real service replicas with latency + self-ID"],
     ["Streaming SSR must not be buffered", "Pipe a Node stream directly to the response"],
     ["SSG selectable only when cache warm", "Pre-build static pages on server startup"],
     ["Docker Compose plugin missing on host", "Verified via docker build/run; documented one-line install"]],
    col_w=[5.2, 6.4], top=1.5, fsize=14, hfsize=15)

# ---------------------------------------------------------------- slide 13 ---
s = content_slide("Remaining Work & Timeline", 13)
bullets(s, [
    "Install Docker Compose; run the full origin + edge + proxy + Redis stack.  (Week 9–10)",
    "Network-condition & device experiments — 2G/3G/4G, mobile/desktop.  (Week 11–12)",
    "Load testing with Apache Bench; cold-vs-warm cache experiments.  (Week 12–13)",
    "Adaptive-vs-static comparative benchmarking and statistical analysis.  (Week 13–14)",
    "Graphs, results interpretation and final report.  (Week 14–16)",
], size=20, gap=14)

# ---------------------------------------------------------------- slide 14 ---
s = content_slide("Conclusion", 14)
bullets(s, [
    "The ARE has moved from design to a working, locally verifiable runtime.",
    "It analyses request context and selects one of six strategies via a pure decision engine.",
    "Functional correctness shown by live strategy switching and passing unit tests.",
    "Preliminary metrics confirm the expected performance trade-offs.",
    "Foundation is solid; remaining work (full experimentation & analysis) is on schedule.",
], size=21, gap=14)

# ---------------------------------------------------------------- slide 15 ---
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, PRIMARY)
rect(s, Inches(4.5), Inches(3.95), Inches(4.33), Inches(0.06), ACCENT)
tb, tf = textbox(s, Inches(1), Inches(2.7), Inches(11.3), Inches(2))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Thank You"; _set(r, 48, True, WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "Questions & Discussion"; _set(r, 22, color=LIGHT, italic=True)

prs.save(OUT)
print("Saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
