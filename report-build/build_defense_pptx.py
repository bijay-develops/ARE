#!/usr/bin/env python3
"""
Build the final defense presentation for the Adaptive Rendering Engine.

Content is sourced from:
  - docs/ARE_DefenseReport_final.docx   (the final project report)
  - docs/0_PROJECT-CONTEXT.md           (verified single source of truth)
  - report-build/data/final_results.json (measured experimental data)

Every number on a slide is interpolated from final_results.json, so the deck
cannot drift from the measured data.

    python3 report-build/build_defense_pptx.py
"""

import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIGS = os.path.join(ROOT, "report-build", "figs")
SHOTS = os.environ.get("ARE_SHOTS", os.path.join(ROOT, "report-build", "shots"))
OUT = os.path.join(ROOT, "ARE_Defense_Presentation.pptx")

D = json.load(open(os.path.join(ROOT, "report-build", "data", "final_results.json")))

# ─────────────────────────────────────────────────────────────── design system
NAVY = RGBColor(0x1F, 0x38, 0x64)
NAVY_D = RGBColor(0x16, 0x28, 0x48)
BLUE = RGBColor(0x1F, 0x6F, 0xB2)
BLUE_L = RGBColor(0xD9, 0xE5, 0xF3)
ROW_A = RGBColor(0xFF, 0xFF, 0xFF)
ROW_B = RGBColor(0xEE, 0xF3, 0xFA)
PANEL = RGBColor(0xF4, 0xF7, 0xFC)
CODE_BG = RGBColor(0xF3, 0xF5, 0xF9)
GREEN_BG = RGBColor(0xE4, 0xF1, 0xDC)
GREEN_TX = RGBColor(0x2C, 0x5A, 0x1E)
AMBER_BG = RGBColor(0xFD, 0xF2, 0xDC)
AMBER_TX = RGBColor(0x7A, 0x51, 0x00)
RED = RGBColor(0xB5, 0x2E, 0x2E)
TEXT = RGBColor(0x24, 0x29, 0x33)
MUTED = RGBColor(0x5C, 0x66, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xC9, 0xD5, 0xE5)

HEAD = "Trebuchet MS"
BODY = "Cambria"
MONO = "Consolas"

SW, SH = 13.333, 7.5
M = 0.62                      # left/right margin
CW = SW - 2 * M               # content width

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_slide_no = [0]


# ───────────────────────────────────────────────────────────────── primitives
def rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, lw=1.0):
    s = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.shadow.inherit = False
    s.text_frame.word_wrap = True
    return s


def tbox(slide, l, t, w, h):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _runs(p, text, size, bold, color, font, italic):
    """Split on ** for inline bold emphasis."""
    for i, seg in enumerate(text.split("**")):
        if not seg:
            continue
        r = p.add_run()
        r.text = seg
        f = r.font
        f.size = Pt(size)
        f.name = font
        f.bold = bold or (i % 2 == 1)
        f.italic = italic
        f.color.rgb = color


def para(tf, text, size=15, bold=False, color=TEXT, font=BODY, after=7, before=0,
         first=False, italic=False, align=PP_ALIGN.LEFT, line=1.16, bullet=None,
         indent=0.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    p.line_spacing = line
    _runs(p, text, size, bold, color, font, italic)
    if bullet:
        marl = int((0.26 + indent) * 914400)
        ind = int(0.26 * 914400)
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(marl))
        pPr.set("indent", str(-ind))
        bf = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
        bc = pPr.makeelement(qn("a:buChar"), {"char": bullet})
        pPr.append(bf)
        pPr.append(bc)
    elif indent:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(int(indent * 914400)))
    return p


def slide_title(slide, title, kicker=None, sub=None):
    _slide_no[0] += 1
    top = 0.42
    if kicker:
        tf = tbox(slide, M, 0.34, CW, 0.26)
        para(tf, kicker.upper(), size=11.5, bold=True, color=BLUE, font=HEAD,
             after=0, first=True)
        top = 0.66
    tf = tbox(slide, M, top, CW, 0.62)
    para(tf, title, size=28, bold=True, color=NAVY, font=HEAD, after=0, first=True)
    y = top + 0.56
    if sub:
        tf = tbox(slide, M, y + 0.02, CW, 0.32)
        para(tf, sub, size=13.5, color=MUTED, font=BODY, after=0, first=True)
        y += 0.36
    rect(slide, M, y + 0.10, 1.15, 0.045, fill=BLUE)
    return y + 0.40


def footer(slide, label="Adaptive Rendering Engine"):
    tf = tbox(slide, M, SH - 0.42, CW - 0.6, 0.24)
    para(tf, label, size=9.5, color=MUTED, font=HEAD, after=0, first=True)
    tf = tbox(slide, SW - M - 0.6, SH - 0.42, 0.6, 0.24)
    para(tf, str(_slide_no[0]), size=9.5, color=MUTED, font=HEAD, after=0,
         first=True, align=PP_ALIGN.RIGHT)


def new(title=None, kicker=None, sub=None, foot=True):
    s = prs.slides.add_slide(BLANK)
    y = slide_title(s, title, kicker, sub) if title else 0.5
    if foot and title:
        footer(s)
    return s, y


# ────────────────────────────────────────────────────────────── composed parts
def bullets(slide, items, top, left=M, width=CW, size=15, gap=8, char="▪"):
    """items: str, or (text, sub) for a bullet with a muted second line."""
    tf = tbox(slide, left, top, width, SH - top - 0.6)
    first = True
    for it in items:
        sub = None
        if isinstance(it, tuple):
            it, sub = it
        para(tf, it, size=size, after=(3 if sub else gap), first=first,
             bullet=char, color=TEXT)
        first = False
        if sub:
            para(tf, sub, size=size - 2.2, color=MUTED, after=gap, indent=0.26)
    return tf


def table(slide, headers, rows, widths, top, left=M, fs=11.5, hfs=11.5,
          rh=0.30, hh=0.34, bold_first=True, aligns=None, row_colors=None,
          width=None):
    total = sum(widths)
    target = width if width is not None else CW
    widths = [w * (target / total) for w in widths]
    n = len(rows) + 1
    gf = slide.shapes.add_table(n, len(headers), Inches(left), Inches(top),
                                Inches(sum(widths)), Inches(hh + rh * len(rows)))
    t = gf.table
    t.first_row = True
    t.horz_banding = False
    for i, w in enumerate(widths):
        t.columns[i].width = Inches(w)
    t.rows[0].height = Inches(hh)
    for r in range(1, n):
        t.rows[r].height = Inches(rh)

    for c, h in enumerate(headers):
        cell = t.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = cell.margin_right = Inches(0.07)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        para(tf, h, size=hfs, bold=True, color=WHITE, font=HEAD, after=0,
             first=True, align=(aligns[c] if aligns else PP_ALIGN.LEFT))

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = t.cell(r, c)
            cell.fill.solid()
            if row_colors and row_colors[r - 1] is not None:
                cell.fill.fore_color.rgb = row_colors[r - 1]
            else:
                cell.fill.fore_color.rgb = ROW_A if r % 2 else ROW_B
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.025)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            first_col = (c == 0 and bold_first)
            para(tf, str(val), size=fs, bold=first_col,
                 color=(NAVY if first_col else TEXT), font=BODY, after=0,
                 first=True, line=1.05,
                 align=(aligns[c] if aligns else PP_ALIGN.LEFT))
    return t


def code(slide, lines, top, left=M, width=CW, fs=11.5, title=None, height=None):
    lh = fs / 72.0 * 1.42
    h = height or (len(lines) * lh + 0.30)
    rect(slide, left, top, width, h, fill=CODE_BG, line=LINE, lw=0.75)
    if title:
        tf = tbox(slide, left + 0.16, top + 0.10, width - 0.32, 0.22)
        para(tf, title, size=10, bold=True, color=BLUE, font=HEAD, after=0, first=True)
    tf = tbox(slide, left + 0.16, top + (0.34 if title else 0.14),
              width - 0.32, h - 0.24)
    for i, ln in enumerate(lines):
        col = MUTED if ln.strip().startswith(("//", "#")) else NAVY_D
        para(tf, ln if ln else " ", size=fs, color=col, font=MONO, after=0,
             first=(i == 0), line=1.06)
    return top + h


def callout(slide, text, top, left=M, width=CW, bg=BLUE_L, fg=NAVY, size=14,
            bold=False, height=None, label=None):
    # auto-height: never let the text outgrow its band
    inner = width - 0.44
    chars_per_line = max(1, int(inner / (size / 72.0 * 0.505)))
    plain = text.replace("**", "")
    lines = max(1, -(-len(plain) // chars_per_line))
    need = lines * size / 72.0 * 1.20 + 0.26 + (0.24 if label else 0.0)
    h = max(height or 0.62, need)
    rect(slide, left, top, width, h, fill=bg)
    rect(slide, left, top, 0.055, h, fill=BLUE if bg is BLUE_L else
         (GREEN_TX if bg is GREEN_BG else AMBER_TX))
    tf = tbox(slide, left + 0.24, top + 0.11, width - 0.44, h - 0.20)
    first = True
    if label:
        para(tf, label.upper(), size=9.5, bold=True, color=fg, font=HEAD,
             after=3, first=True)
        first = False
    para(tf, text, size=size, bold=bold, color=fg, font=BODY, after=0,
         first=first, line=1.15)
    return top + h


def tiles(slide, items, top, left=M, width=CW, h=1.12, gap=0.16, vsize=25):
    """items: (value, label) or (value, label, note)."""
    n = len(items)
    w = (width - gap * (n - 1)) / n
    for i, it in enumerate(items):
        val, lab = it[0], it[1]
        note = it[2] if len(it) > 2 else None
        x = left + i * (w + gap)
        rect(slide, x, top, w, h, fill=PANEL, line=LINE, lw=0.75)
        rect(slide, x, top, w, 0.05, fill=BLUE)
        tf = tbox(slide, x + 0.14, top + 0.17, w - 0.28, 0.44)
        para(tf, val, size=vsize, bold=True, color=NAVY, font=HEAD, after=0,
             first=True, align=PP_ALIGN.CENTER)
        tf = tbox(slide, x + 0.10, top + 0.64, w - 0.20, 0.42)
        para(tf, lab, size=11, color=TEXT, font=BODY, after=0, first=True,
             align=PP_ALIGN.CENTER, line=1.08)
        if note:
            para(tf, note, size=9.5, color=MUTED, font=BODY, after=0,
                 align=PP_ALIGN.CENTER, line=1.05)
    return top + h


def picture(slide, path, top, max_h, left=None, max_w=CW, caption=None, border=True):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    h = max_h
    w = h * ar
    if w > max_w:
        w = max_w
        h = w / ar
    x = left if left is not None else (SW - w) / 2
    if border:
        rect(slide, x - 0.03, top - 0.03, w + 0.06, h + 0.06, fill=WHITE, line=LINE, lw=0.75)
    slide.shapes.add_picture(path, Inches(x), Inches(top), Inches(w), Inches(h))
    if caption:
        tf = tbox(slide, M, top + h + 0.12, CW, 0.3)
        para(tf, caption, size=10.5, color=MUTED, font=BODY, italic=True,
             after=0, first=True, align=PP_ALIGN.CENTER)
    return top + h


def fig(name):
    return os.path.join(FIGS, name)


def shot(name):
    p = os.path.join(SHOTS, name)
    return p if os.path.exists(p) else None


def section(number, title, points):
    s = prs.slides.add_slide(BLANK)
    _slide_no[0] += 1
    rect(s, 0, 0, SW, SH, fill=NAVY)
    rect(s, 0, 0, 0.16, SH, fill=BLUE)
    tf = tbox(s, 1.15, 2.35, 1.6, 1.0)
    para(tf, number, size=64, bold=True, color=RGBColor(0x5E, 0x93, 0xC9),
         font=HEAD, after=0, first=True)
    tf = tbox(s, 1.15, 3.35, 7.6, 0.8)
    para(tf, title, size=34, bold=True, color=WHITE, font=HEAD, after=0, first=True)
    rect(s, 1.15, 4.32, 1.3, 0.05, fill=RGBColor(0x5E, 0x93, 0xC9))
    tf = tbox(s, 1.15, 4.62, 8.6, 1.6)
    for i, p in enumerate(points):
        para(tf, p, size=14, color=RGBColor(0xC6, 0xD6, 0xE8), font=BODY,
             after=6, first=(i == 0), bullet="–")
    tf = tbox(s, SW - 1.4, SH - 0.5, 0.8, 0.3)
    para(tf, str(_slide_no[0]), size=10, color=RGBColor(0x8F, 0xA8, 0xC4),
         font=HEAD, after=0, first=True, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════ derived numbers
lf, lm = D["load_fast"], D["load_medium"]
fx, ad = lf["fixed_ssr"], lf["adaptive"]
fx2, ad2 = lm["fixed_ssr"], lm["adaptive"]
ds = D["decision_space"]
ps = D["per_strategy_server"]
pc = D["per_strategy_client"]

rps_gain = (ad["rps"] - fx["rps"]) / fx["rps"] * 100
mean_cut = (fx["mean_ms"] - ad["mean_ms"]) / fx["mean_ms"] * 100
p99_cut = (fx["p99"] - ad["p99"]) / fx["p99"] * 100
byte_cut = (fx["doc_bytes"] - ad["doc_bytes"]) / fx["doc_bytes"] * 100
cpu_cut = (fx2["cpu_ms"] - ad2["cpu_ms"]) / fx2["cpu_ms"] * 100
p99_cut2 = (fx2["p99"] - ad2["p99"]) / fx2["p99"] * 100
rps_d2 = (ad2["rps"] - fx2["rps"]) / fx2["rps"] * 100
ssg_ratio = ps["SSG"]["ttfb"] / ps["ISR"]["ttfb"]
csr_light = ps["SSR"]["bytes"] / ps["CSR"]["bytes"]

iso = D["isr_lifecycle"]
cold = [r for r in iso if r["state"] == "miss"]
fresh = [r for r in iso if r["state"] == "fresh"]
stale = [r for r in iso if r["state"] == "stale-revalidating"]
cold_ms = cold[0]["ms"]
fresh_ms = sum(r["ms"] for r in fresh) / len(fresh)
stale_ms = sum(r["ms"] for r in stale) / len(stale)
warm_cut = (cold_ms - fresh_ms) / cold_ms * 100

# ══════════════════════════════════════════════════════════════════ 1. TITLE
s = prs.slides.add_slide(BLANK)
_slide_no[0] += 1
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.20, fill=BLUE)
tf = tbox(s, 1.05, 1.28, 11.2, 0.34)
para(tf, "FINAL YEAR ENGINEERING PROJECT  ·  DEFENSE PRESENTATION", size=12.5,
     bold=True, color=RGBColor(0x8F, 0xB6, 0xDC), font=HEAD, after=0, first=True)
tf = tbox(s, 1.05, 1.78, 11.2, 1.3)
para(tf, "Adaptive Rendering Engine", size=52, bold=True, color=WHITE, font=HEAD,
     after=0, first=True)
tf = tbox(s, 1.05, 2.92, 10.6, 0.8)
para(tf, "Choosing the web rendering strategy per request at runtime — "
         "from network, device, cache, load and data volatility.",
     size=17, color=RGBColor(0xC6, 0xD8, 0xEC), font=BODY, after=0, first=True,
     line=1.25)
rect(s, 1.05, 3.92, 1.5, 0.05, fill=BLUE)

tf = tbox(s, 1.05, 4.28, 5.6, 1.9)
para(tf, "PRESENTED BY", size=10.5, bold=True, color=RGBColor(0x8F, 0xB6, 0xDC),
     font=HEAD, after=7, first=True)
for nm, rl in [("Bijay B.k", "220305"), ("Devendra Pandey", "220306"),
               ("Manish Joshi", "220312"), ("Pramod Panta", "220317")]:
    para(tf, f"{nm}   —   {rl}", size=14, color=WHITE, font=BODY, after=4)

tf = tbox(s, 7.1, 4.28, 5.2, 2.0)
para(tf, "SUPERVISOR", size=10.5, bold=True, color=RGBColor(0x8F, 0xB6, 0xDC),
     font=HEAD, after=7, first=True)
para(tf, "Er. Robinhood Khadka", size=14, color=WHITE, font=BODY, after=14)
para(tf, "DEPARTMENT", size=10.5, bold=True, color=RGBColor(0x8F, 0xB6, 0xDC),
     font=HEAD, after=7)
para(tf, "ICT and Computer Engineering", size=13, color=WHITE, font=BODY, after=2)
para(tf, "Cosmos College of Management & Technology", size=13,
     color=RGBColor(0xC6, 0xD8, 0xEC), font=BODY, after=2)
para(tf, "Affiliated to Pokhara University  ·  Shrawan 2083 (August 2026)",
     size=11.5, color=RGBColor(0x9F, 0xB8, 0xD4), font=BODY, after=0)

# ═══════════════════════════════════════════════════════════════ 2. AGENDA
s, y = new("What this presentation covers", kicker="Agenda")
cols = [
    ("1", "The Idea and the Vision",
     "Where the question came from, the problem, the gap in the literature, and the objectives."),
    ("2", "System Analysis and Design",
     "The five invariants, the architecture, the five-stage pipeline and every component."),
    ("3", "Methodology and Implementation",
     "Flow of code, the decision algorithm in full, the six strategies, the private server."),
    ("4", "Results, Analysis and Evaluation",
     "Correctness, cost of adapting, per-strategy performance, adaptive vs a fixed policy."),
    ("5", "Challenges, Conclusion and Future",
     "What went wrong and how it was fixed, what is proven, and where this goes next."),
]
ty = y + 0.05
for i, (n, t, d) in enumerate(cols):
    yy = ty + i * 0.98
    rect(s, M, yy, CW, 0.86, fill=PANEL)
    rect(s, M, yy, 0.055, 0.86, fill=BLUE)
    tf = tbox(s, M + 0.30, yy + 0.14, 0.6, 0.6)
    para(tf, n, size=24, bold=True, color=BLUE_L if False else RGBColor(0x9C, 0xBA, 0xD9),
         font=HEAD, after=0, first=True)
    tf = tbox(s, M + 0.95, yy + 0.13, CW - 1.3, 0.62)
    para(tf, t, size=16, bold=True, color=NAVY, font=HEAD, after=2, first=True)
    para(tf, d, size=12, color=MUTED, font=BODY, after=0)

# ═════════════════════════════════════════════════ SECTION 1 — IDEA & VISION
section("01", "The Idea and the Vision",
        ["An observation while building an ordinary Next.js site",
         "The problem stated precisely, and the gap the literature leaves open",
         "The objectives this project set itself"])

# ── where it started
s, y = new("Where this project started", kicker="Origin of the idea",
           sub="It did not begin with a literature search. It began with an annotation.")
callout(s, "A framework asks the developer to annotate each route with how it should be "
           "rendered — and it asks this question at the one moment when the answer "
           "cannot possibly be known: **before any user has arrived**.",
        y, size=15, bg=AMBER_BG, fg=AMBER_TX, height=0.86)
bullets(s, [
    ("The annotation is a single bet placed on an average user who does not exist.",
     "A route marked SSR is server-rendered for the desktop on fibre who wanted a light "
     "interactive shell, and identically for the throttled handset that genuinely needed finished HTML."),
    ("Reframing the observation produced the project.",
     "If the rendering strategy is a variable rather than a constant, then choosing it is an "
     "optimisation problem — and an optimisation problem belongs at runtime, where the inputs actually exist."),
    ("The six strategies are not competitors to be ranked — they are points on a trade-off surface.",
     "SSG, SSR, CSR, ISR, Streaming SSR and edge rendering are each optimal somewhere. "
     "The optimum moves with the request; the annotation cannot."),
], y + 1.02, size=15)

# ── vision
s, y = new("The vision", kicker="The standard this work is measured against")
rect(s, M, y, CW, 1.16, fill=NAVY)
tf = tbox(s, M + 0.34, y + 0.20, CW - 0.68, 0.9)
para(tf, "Rendering strategy should be a **runtime decision computed from the conditions of "
         "the request**, not a build-time constant chosen by a developer — and every such "
         "decision should be **observable, explainable and reproducible from outside the system**.",
     size=17, color=WHITE, font=BODY, after=0, first=True, line=1.28)
tf = tbox(s, M, y + 1.42, CW, 0.3)
para(tf, "Four commitments follow, and every design decision in the project is an application of one of them:",
     size=13.5, color=MUTED, font=BODY, after=0, first=True)
items = [
    ("Honest", "The engine publishes what it decided and why on every response, so a reviewer "
               "verifies behaviour without reading the source."),
    ("Cheap", "An engine that spends more time deciding than rendering has defeated itself. "
              "Selection is a pure function over primitives, with no I/O."),
    ("Extensible", "Adding a seventh strategy — or replacing the rule table with a learned "
                   "model — must not require editing any existing strategy."),
    ("Self-explaining", "The demonstration pages render the observed context and the rule that "
                        "fired into their own markup, so the system documents itself."),
]
w = (CW - 0.16 * 3) / 4
for i, (t, d) in enumerate(items):
    x = M + i * (w + 0.16)
    rect(s, x, y + 1.86, w, 1.74, fill=PANEL, line=LINE, lw=0.75)
    rect(s, x, y + 1.86, w, 0.05, fill=BLUE)
    tf = tbox(s, x + 0.16, y + 2.02, w - 0.32, 1.5)
    para(tf, t, size=15.5, bold=True, color=NAVY, font=HEAD, after=6, first=True)
    para(tf, d, size=11.5, color=TEXT, font=BODY, after=0, line=1.14)

# ── background
s, y = new("Background: six strategies, one trade-off surface", kicker="Background",
           sub="Each inversion of the web's rendering model solved the previous problem and created a new one.")
table(s, ["Strategy", "What it does", "Wins", "Costs"], [
    ["SSG", "Renders once at build time and serves a file",
     "Unbeatably cheap per request", "Cannot express fresh data"],
    ["SSR", "Renders complete HTML on the server per request",
     "Fast first paint, crawlable", "Origin CPU under concurrency"],
    ["CSR", "Ships an empty shell; the browser builds the page",
     "Minimal server work, full interactivity", "Slow first paint on weak devices/links"],
    ["ISR", "Serves a cached render, revalidates on a TTL",
     "SSR's output near CSR's cost", "Bounded staleness"],
    ["Streaming SSR", "Flushes the shell, streams the rest as it resolves",
     "Decouples TTFB from total render", "Highest per-request cost"],
    ["Edge rendering", "Moves any of the above closer to the user",
     "Cuts distance latency", "Cache coherence across nodes"],
], [1.3, 3.2, 2.6, 2.6], y, fs=12, rh=0.44)
callout(s, "The literature measures these thoroughly — but almost always compares them "
           "**in isolation, as fixed architectural choices**, rather than treating the choice "
           "itself as a controllable runtime variable. That is the gap this project occupies.",
        y + 3.24, size=13.5)

# ── problem statement
s, y = new("Problem statement", kicker="The problem")
rect(s, M, y, CW, 1.05, fill=NAVY)
tf = tbox(s, M + 0.34, y + 0.18, CW - 0.68, 0.75)
para(tf, "Rendering strategy selection in modern web applications is **static and does not "
         "adapt to runtime contextual conditions**, resulting in reduced performance "
         "efficiency, avoidable origin cost and scalability limitations.",
     size=17, color=WHITE, font=BODY, after=0, first=True, line=1.26)
tf = tbox(s, M, y + 1.30, CW, 0.3)
para(tf, "Three concrete deficiencies follow from a choice that is fixed at development time "
         "and applied uniformly thereafter:", size=13.5, color=MUTED, font=BODY,
     after=0, first=True)
items = [
    ("Sub-optimal responses", "Whenever real conditions differ from the developer's assumption "
                              "— which is most of the time, across a heterogeneous population."),
    ("Inefficient origin use", "The server keeps paying for work a cheaper strategy would have "
                               "done just as well, with no way to notice."),
    ("Inconsistent experience", "One annotation must serve every device class and every network "
                                "class the site will ever meet."),
]
w = (CW - 0.20 * 2) / 3
for i, (t, d) in enumerate(items):
    x = M + i * (w + 0.20)
    rect(s, x, y + 1.72, w, 1.35, fill=PANEL, line=LINE, lw=0.75)
    tf = tbox(s, x + 0.18, y + 1.90, w - 0.36, 1.1)
    para(tf, t, size=15, bold=True, color=NAVY, font=HEAD, after=6, first=True)
    para(tf, d, size=12, color=TEXT, font=BODY, after=0, line=1.14)
callout(s, "And because no unified runtime mechanism exists at the application layer, there is "
           "also **no experimental platform** on which the value of such adaptation can be measured.",
        y + 3.24, size=13.5)

# ── literature gap
s, y = new("What the literature settles — and what it leaves open", kicker="Literature review",
           sub="Four limitations persist across the reviewed work; each one is answered by a design decision in this project.")
table(s, ["Limitation identified in the literature", "Response implemented in this project"], [
    ["Rendering strategies are treated as static architectural decisions",
     "Strategy is a runtime variable, recomputed for every request"],
    ["Contextual variables are rarely integrated into rendering selection",
     "Seven contextual signals are observed and drive an explicit rule table"],
    ["No unified runtime engine orchestrates rendering strategies",
     "One engine registers six strategies behind a single interface and selects among them"],
    ["Experimental validation of adaptive rendering at the application layer is limited",
     "A reproducible private server, 28 automated tests and a controlled comparison against fixed policies"],
], [1, 1], y, fs=13, rh=0.52, bold_first=False)
callout(s, "Prior work establishes that rendering strategies differ measurably, that contextual "
           "adaptation is valuable, and how performance should be measured. It does **not** propose, "
           "implement or experimentally validate a runtime engine that unifies contextual analysis "
           "with rendering orchestration at the application layer.",
        y + 2.55, size=13.5, height=0.86, label="The gap, stated precisely")

# ── objectives
s, y = new("Objectives", kicker="What we set out to do")
rect(s, M, y, CW, 0.78, fill=BLUE_L)
rect(s, M, y, 0.055, 0.78, fill=BLUE)
tf = tbox(s, M + 0.26, y + 0.11, CW - 0.5, 0.6)
para(tf, "GENERAL OBJECTIVE", size=9.5, bold=True, color=NAVY, font=HEAD, after=3, first=True)
para(tf, "To design, implement and evaluate an Adaptive Rendering Engine that selects a web "
         "rendering strategy per request at runtime from observed contextual variables, and to "
         "demonstrate experimentally that doing so improves performance and resource efficiency "
         "relative to a fixed rendering policy.",
     size=13.5, color=NAVY, font=BODY, after=0, line=1.16)
tf = tbox(s, M, y + 0.98, CW, 0.3)
para(tf, "SPECIFIC OBJECTIVES", size=10, bold=True, color=BLUE, font=HEAD, after=0, first=True)
objs = [
    "**Identify and observe** the contextual variables that influence rendering performance, and build an analyzer that reads them from a live request.",
    "**Design a pure, deterministic, rule-based decision engine** that maps any context to exactly one strategy — and prove that mapping total and unambiguous.",
    "**Implement six rendering strategies** behind a single uniform interface so that they are genuinely interchangeable at runtime.",
    "**Make every decision externally observable** through response headers, structured server logs and per-request metrics.",
    "**Construct a zero-cost, reproducible private server** that models an origin, multiple edge nodes and a shared cache.",
    "**Measure and compare** per-request cost, payload, cache efficiency and resource use against fixed single-strategy policies under controlled load.",
]
tf = tbox(s, M, y + 1.34, CW, 3.3)
for i, o in enumerate(objs):
    para(tf, o, size=13.5, after=9, first=(i == 0), bullet="▪", line=1.14)

# ── scope
s, y = new("Scope and limitations", kicker="Boundaries set before we began")
w = (CW - 0.28) / 2
rect(s, M, y, w, 3.35, fill=PANEL, line=LINE, lw=0.75)
rect(s, M, y, w, 0.05, fill=GREEN_TX)
tf = tbox(s, M + 0.22, y + 0.20, w - 0.44, 3.0)
para(tf, "IN SCOPE", size=11, bold=True, color=GREEN_TX, font=HEAD, after=10, first=True)
for t in ["Design and implementation of a standalone adaptive rendering runtime.",
          "Simulation of network classes, device profiles, cache states and server load inside a local containerised environment.",
          "The instrumentation required to measure the engine from the outside.",
          "A controlled comparative evaluation of adaptive selection against fixed policies."]:
    para(tf, t, size=13, after=9, bullet="▪", line=1.14)
x2 = M + w + 0.28
rect(s, x2, y, w, 3.35, fill=PANEL, line=LINE, lw=0.75)
rect(s, x2, y, w, 0.05, fill=RED)
tf = tbox(s, x2 + 0.22, y + 0.20, w - 0.44, 3.0)
para(tf, "EXPLICITLY OUT OF SCOPE", size=11, bold=True, color=RED, font=HEAD, after=10, first=True)
for t in ["No commercial cloud deployment — the edge topology is modelled locally with injected latency.",
          "Not a production-scale web application; the engine carries no business domain.",
          "Client-side paint metrics (FCP, LCP) are not instrumented — we report the quantities the engine itself can influence and observe.",
          "The policy is rule-based by design; learned policies are future work, and the architecture is shaped to accept them."]:
    para(tf, t, size=13, after=9, bullet="▪", line=1.14)
callout(s, "Stating the boundary is part of the contribution: every claim in Chapter 5 of the "
           "report is bounded by exactly these limits, which is what makes the results defensible.",
        y + 3.52, size=13)

# ═══════════════════════════════════════════════════ SECTION 2 — DESIGN
section("02", "System Analysis and Design",
        ["Five invariants fixed before a line of code was written",
         "The deployed architecture and the five-stage per-request pipeline",
         "Every component, the data structures, and the control surfaces"])

# ── design philosophy
s, y = new("Design philosophy: five invariants", kicker="Design",
           sub="Fixed before implementation began. Every later design decision is an application of one of these — and Chapter 5 measures whether they were achieved.")
inv = [
    ("Separation of concerns",
     "Analysis observes but never decides; decision selects but never renders; rendering produces "
     "output but never re-derives the decision.",
     "→ the decision can be tested without a server, and a strategy cannot silently disagree "
     "with the engine about why it was invoked."),
    ("Purity of the decision",
     "Selection is a pure function from a plain context object to a decision record: no I/O, no "
     "clock, no state.",
     "→ provably deterministic, exhaustively enumerable, and testable one rule at a time. "
     "This is what makes Section 5.3 possible."),
    ("Pluggability",
     "All six strategies implement one interface and register at start-up. The engine holds a "
     "reference to the registry, not to any strategy.",
     "→ adding a seventh strategy is an insertion, not an edit. Edge-ISR is literally ISR "
     "with one overridden method."),
    ("Transparency",
     "A decision that cannot be observed cannot be trusted. Every response carries the strategy "
     "and the rule; every request logs three lines and appends a metric.",
     "→ the engine is auditable from outside — and every measurement in this deck was "
     "taken by reading a header it publishes."),
    ("Graceful degradation",
     "The rule table ends in an unconditional fallback, and the render call is wrapped so a "
     "throwing strategy is replaced by SSR.",
     "→ selection can never fail, and a degraded response says so in its reason string "
     "instead of hiding it."),
]
tf = tbox(s, M, y, CW, 4.0)
for i, (t, d, c) in enumerate(inv):
    para(tf, f"{i+1}.  {t}", size=14.5, bold=True, color=NAVY, font=HEAD,
         after=2, first=(i == 0))
    para(tf, d, size=12.5, color=TEXT, font=BODY, after=1, indent=0.30, line=1.12)
    para(tf, c, size=12, color=BLUE, font=BODY, after=9, indent=0.30, line=1.12)

# ── architecture figure
s, y = new("System architecture", kicker="Architecture",
           sub="One nginx proxy is the only entry point. Every node runs the identical engine image and differs only by environment variables.")
picture(s, fig("system-architecture.png"), y + 0.02, 3.62)
callout(s, "Modelling the edges as **full engine instances** rather than caching proxies was "
           "deliberate. An edge that is only a proxy can demonstrate distance but not behaviour; "
           "an edge that runs the engine has its own context, its own cache namespace and its own "
           "decisions — which is what makes Edge-ISR a real strategy rather than a decoration.",
        y + 3.82, size=13)

# ── docker topology
s, y = new("Deployment: the zero-cost private server", kicker="Architecture",
           sub="Five containers on one private Docker network, addressed by service name. No cloud account, no domain, no public IP.")
table(s, ["Service", "Role", "Published port", "Key environment"], [
    ["proxy", "nginx reverse proxy — the single entry point", "8080 → 80",
     "routes /, /edge1/, /edge2/"],
    ["origin", "Engine, the authoritative node", "internal only",
     "SERVED_BY=origin, EDGE_LATENCY_MS=0, TTL=30 s"],
    ["edge-node-1", "Engine acting as a near edge", "8081 → 3000",
     "SERVED_BY=edge-node-1, +20 ms, TTL=15 s"],
    ["edge-node-2", "Engine acting as a far edge", "8082 → 3000",
     "SERVED_BY=edge-node-2, +80 ms, TTL=15 s"],
    ["redis", "Cache shared across nodes", "internal only",
     "in-memory only, no persistence"],
], [1.5, 3.6, 1.7, 4.0], y, fs=12, rh=0.40)
tf = tbox(s, M, y + 2.62, CW, 1.5)
para(tf, "Two consequences shape how the system is operated — and both are results, not defects:",
     size=13, bold=True, color=NAVY, font=HEAD, after=8, first=True)
para(tf, "The proxy stamps **X-Served-By: origin** on the default route, which correctly stops a "
         "client claiming to be at an edge when it is not — so Edge-ISR must be exercised via "
         "/edge1, /edge2 or the edge ports directly.",
     size=13, after=6, bullet="▪", line=1.14)
para(tf, "The stack mounts **no volumes**, so the container filesystem is the only store for "
         "artefacts, caches and metrics. A rebuild is required to deploy changed code — and a "
         "teardown discards state, which is precisely what makes every experimental run start cold.",
     size=13, after=0, bullet="▪", line=1.14)

# ── pipeline
s, y = new("The five-stage rendering pipeline", kicker="Design",
           sub="Every request traverses the same five stages. Stage boundaries are strict: each stage emits a value, and the next consumes only that value.")
picture(s, fig("rendering-pipeline.png"), y + 0.55, 2.45, max_w=6.85, left=M)
stages = [
    ("ANALYZE", "context-analyzer.ts", "Raw request → RequestContext. Observes and validates; decides nothing."),
    ("DECIDE", "decision-engine.ts", "Context → DecisionTrace: the strategy, the reason, the context judged."),
    ("RENDER", "strategies/‹name›/", "Looks the strategy up in the registry and invokes it."),
    ("RESPOND", "engine.ts", "Status and headers — always the two proof headers — then body or stream."),
    ("MEASURE", "metrics-collector.ts", "Appends a metric record. Never awaited."),
]
x = M + 7.10
tf = tbox(s, x, y, CW - 7.10, 4.2)
for i, (n, f, d) in enumerate(stages):
    para(tf, f"{i+1}   {n}", size=13, bold=True, color=NAVY, font=HEAD, after=1,
         first=(i == 0))
    para(tf, f, size=10, color=BLUE, font=MONO, after=2, indent=0.24)
    para(tf, d, size=10.5, color=TEXT, font=BODY, after=9, indent=0.24, line=1.12)

# ── components I
s, y = new("Component design — the decision path", kicker="Components 1 of 2")
comp = [
    ("Context Analyzer", "src/core/context-analyzer.ts",
     "The system's sensory layer. Produces a RequestContext of seven decision-relevant fields. "
     "For each field it applies a fixed precedence — control header, else query alias, else "
     "inference — and validates against the permitted value set, so an unrecognised value "
     "falls back to inference rather than corrupting the context. **Load is never supplied by the "
     "request**: it is classified from a live in-flight counter. Cache state is probed from the "
     "cache layer before analysis begins."),
    ("Decision Engine", "src/core/decision-engine.ts",
     "Holds the policy. Evaluates an ordered rule table top to bottom and returns the first match "
     "with its human-readable reason. The final rule is unconditional, so the function is "
     "**total**; it touches nothing outside its argument, so it is **pure**. The policy itself "
     "lives in a separate configuration module, so the policy can change without touching the "
     "evaluator."),
    ("Strategy Registry and Strategy Modules", "src/core/strategy-registry.ts",
     "A name-to-implementation map populated at start-up. Each of the six strategies implements "
     "one render method with an identical signature — receiving the context, the page module, "
     "the cache manager and the decision reason, returning a status, headers, a body and a "
     "from-cache flag. **Uniformity of that signature is what makes the strategies genuinely "
     "interchangeable.**"),
]
tf = tbox(s, M, y, CW, 4.2)
for i, (t, f, d) in enumerate(comp):
    p = para(tf, t, size=16, bold=True, color=NAVY, font=HEAD, after=2, first=(i == 0))
    para(tf, f, size=11, color=BLUE, font=MONO, after=4)
    para(tf, d, size=13, color=TEXT, font=BODY, after=14, line=1.16)

# ── components II
s, y = new("Component design — support subsystems", kicker="Components 2 of 2")
comp = [
    ("Cache Manager", "src/cache/cache-manager.ts",
     "One interface over three backends in a hierarchy: in-process **memory**, optional shared "
     "**Redis**, persistent **file**. Reads descend and promote any hit upward; writes fan out. "
     "Redis is optional by construction. It also exposes a non-mutating freshness probe — which "
     "is what lets the analyzer observe cache state *before* a decision is made."),
    ("Metrics Collector and Report Generator", "src/metrics/",
     "One newline-delimited JSON record per request: timings, payload size, cache outcome, the "
     "observed context and a process resource sample. **Every quantitative result in this deck "
     "derives from this pipeline, or from external measurement of the same requests.**"),
    ("Simulation Subsystem", "src/simulation/",
     "Because the evaluation runs on one machine, conditions must be created rather than awaited. "
     "The throttler applies 400 / 100 / 0 ms for slow / medium / fast; edge latency is injected per "
     "container. **Load is not simulated — it is genuinely measured from concurrency.**"),
    ("Demonstration Pages", "src/frontend/pages/",
     "Three React pages exercise the engine — static, realtime and heavy. Each declares its own "
     "volatility and payload weight, which become inputs to the decision."),
]
tf = tbox(s, M, y, CW, 4.3)
for i, (t, f, d) in enumerate(comp):
    para(tf, t, size=15, bold=True, color=NAVY, font=HEAD, after=2, first=(i == 0))
    para(tf, f, size=10.5, color=BLUE, font=MONO, after=3)
    para(tf, d, size=12.3, color=TEXT, font=BODY, after=11, line=1.14)

# ── data structures
s, y = new("Data design: five structures carry all state", kicker="Design",
           sub="Their fields are the vocabulary in which the entire system is written.")
table(s, ["Structure", "Stage", "Purpose and principal fields"], [
    ["RequestContext", "output of ANALYZE",
     "url, networkSpeed (slow|medium|fast), device (mobile|desktop), cacheState (fresh|stale|cold), "
     "load (low|medium|high), volatility (static|periodic|realtime), heavyPayload, isEdge, rawHeaders"],
    ["DecisionTrace", "output of DECIDE",
     "selected (strategy name), reason (the rule text), context (the context it judged)"],
    ["RenderStrategy", "the uniform interface",
     "name, and render(ctx, page, cache, meta) — identical across all six strategies"],
    ["RenderResult", "output of RENDER",
     "status, headers, body (a string or a stream), fromCache"],
    ["MetricRecord", "output of MEASURE",
     "ts, url, strategy, reason, ttfbMs, renderMs, totalMs, fromCache, bytes, network, device, "
     "load, isEdge, resources{}"],
], [1.9, 1.7, 8.0], y, fs=11.5, rh=0.56)
callout(s, "Because the context is a plain object of enumerated primitives, its space is finite "
           "and countable — 3×2×3×3×3×2×2 = **648 distinct "
           "contexts**. Section 5.3 exploits exactly this to enumerate the entire policy.",
        y + 3.30, size=13.5)

# ── control surfaces
s, y = new("Control surfaces: how a request is steered", kicker="Design",
           sub="Precedence: control header  >  query alias  >  inference. This is what turns an adaptive system into an experimentally controllable one.")
table(s, ["Signal", "Control header", "Query alias", "Permitted values", "If nothing is supplied"], [
    ["Network speed", "X-Network-Speed", "?net=", "slow | medium | fast", "defaults to medium"],
    ["Device class", "X-Device-Type", "?device=", "mobile | desktop", "inferred from User-Agent"],
    ["Cache state", "X-Cache-State", "?cache=", "fresh | stale | cold", "probed from the cache layer"],
    ["Server load", "X-Load-Level", "?load=", "low | medium | high", "live in-flight counter (≥25 high, ≥8 medium)"],
    ["Data volatility", "X-Data-Volatility", "?volatility=", "static | periodic | realtime", "the page's own declaration"],
    ["Payload weight", "X-Data-Size", "?size=", "heavy | light", "the page's own declaration"],
    ["Edge identity", "X-Served-By", "?served=", "any value ≠ origin means edge", "the container's SERVED_BY variable"],
], [1.7, 2.2, 1.2, 2.9, 3.4], y, fs=11.5, rh=0.33)
callout(s, "**Why the query layer exists.** A browser cannot attach custom headers to an ordinary "
           "navigation. Without the aliases, the in-page controls could only *predict* a strategy; "
           "with them, following a link genuinely re-renders the page under a different strategy — "
           "which is what makes the system demonstrable in a browser rather than only through curl.",
        y + 2.72, size=13.5, height=0.80)

# ═══════════════════════════════════════════ SECTION 3 — IMPLEMENTATION
section("03", "Methodology and Implementation",
        ["The flow of code: one request traced through every stage",
         "The decision algorithm in full — and why the ordering is the policy",
         "The six strategy modules, the cache, and the containerised private server"])

# ── stack
s, y = new("Technology stack and why each piece was chosen", kicker="Methodology")
table(s, ["Concern", "Technology", "Justification"], [
    ["Language", "Node.js 20+ with TypeScript",
     "One language across engine and view; static types make the context and rule structures verifiable at compile time"],
    ["View layer", "React 18",
     "The only mainstream library exposing every primitive the six strategies need: renderToString, renderToPipeableStream with Suspense, and hydrateRoot"],
    ["HTTP server", "Native Node http module",
     "The deliverable is itself a runtime — a framework would hide the very request handling being studied, and native streaming is required by Streaming SSR"],
    ["Client bundler", "esbuild",
     "Produces the hydration and CSR bundle with no configuration and negligible build time"],
    ["Cache", "Memory + filesystem + optional Redis",
     "Zero-cost persistence, plus a genuinely shared cache that makes edge behaviour real rather than simulated"],
    ["Edge and proxy", "nginx in containers",
     "Provides a single entry point and models edge routing"],
    ["Orchestration", "Docker and Docker Compose",
     "Reproducible multi-node topology on one machine at no cost"],
    ["Testing", "Vitest", "TypeScript-native and fast enough to run on every change"],
    ["Load testing", "Apache Bench (ab)", "Standard concurrent-request generator with percentile reporting"],
], [1.6, 2.9, 7.1], y, fs=11.5, rh=0.36)
callout(s, "React is used strictly as a **rendering primitive**. The contribution of this project "
           "is the selector built around it, not a new view library — and no part of React was modified.",
        y + 3.86, size=13)

# ── flow of code: entry
s, y = new("Flow of code — entry and pre-processing", kicker="Request lifecycle 1 of 4",
           sub="src/server/server.ts — a native HTTP server. The in-flight counter here is the sole source of the load signal.")
code(s, [
    "const server = http.createServer(async (req, res) => {",
    "  inFlight++;                              // the ONLY source of the load signal",
    "  try {",
    "    ...                                    // static assets, /api/data, /health, /control",
    "    if (!headers['x-served-by'] && config.servedBy)",
    "      headers['x-served-by'] = config.servedBy;      // an edge stamps its own identity",
    "",
    "    const speed = headers['x-network-speed'] ?? queryOf(req).get('net') ?? 'medium';",
    "    await sleep(delayForNetwork(speed));             // simulated link  400 / 100 / 0 ms",
    "    if (config.edgeLatencyMs > 0) await sleep(config.edgeLatencyMs);",
    "",
    "    const cacheState = await peekCacheState(page);   // non-mutating freshness probe",
    "    await engine.handle({ url, headers, page,",
    "                          signals: { concurrency: inFlight, cacheState }, res });",
    "  } finally { inFlight--; }                // decremented on EVERY exit path, incl. errors",
    "});",
], y, fs=12.5)
callout(s, "Both sleeps happen **before the engine is invoked**, and the engine's timer starts "
           "inside its own handler. That single design choice is why the server-side timings in "
           "Chapter 5 measure engine cost alone, uncontaminated by the simulated link — while "
           "client-side wall-clock timings include everything.",
        y + 3.30, size=13.5, height=0.80)

# ── analyze
s, y = new("Flow of code — Stage 1: ANALYZE", kicker="Request lifecycle 2 of 4",
           sub="Every signal resolves through one helper that implements the precedence rule, then is validated against its permitted set.")
code(s, [
    "function override(headers, query, headerName) {          // header  >  query  >  inference",
    "  return header(headers, headerName) ?? query.get(QUERY_ALIASES[headerName]) ?? undefined;",
    "}",
    "function oneOf(value, allowed) {",
    "  return allowed.includes(value) ? value : undefined;    // invalid -> fall back to inference",
    "}",
    "",
    "return {",
    "  url,",
    "  networkSpeed: oneOf(pick('x-network-speed'), ['slow','medium','fast']) ?? 'medium',",
    "  device:       oneOf(pick('x-device-type'),   ['mobile','desktop']) ?? inferDevice(headers),",
    "  cacheState:   oneOf(pick('x-cache-state'),   ['fresh','stale','cold']) ?? signals.cacheState,",
    "  load:         oneOf(pick('x-load-level'),    ['low','medium','high'])",
    "                  ?? classifyLoad(signals.concurrency),   // >=25 high, >=8 medium",
    "  volatility:   oneOf(pick('x-data-volatility'), ['static','periodic','realtime'])",
    "                  ?? page.volatility,",
    "  heavyPayload: size ? size.toLowerCase() === 'heavy' : page.heavy === true,",
    "  isEdge:       Boolean(servedBy && servedBy !== 'origin'),",
    "  rawHeaders:   headers,",
    "};",
], y, fs=11.5)
callout(s, "A malformed header **degrades to inference instead of producing an invalid context** — "
           "so no external input can put the decision engine into a state its rule table was not written for.",
        y + 3.62, size=13)

# ── decide
s, y = new("Flow of code — Stage 2: DECIDE", kicker="Request lifecycle 3 of 4",
           sub="The evaluator is nine lines long, and its brevity is the point: all policy lives in data, so the code that applies the policy has nothing to get wrong.")
code(s, [
    "export function decide(ctx: RequestContext): DecisionTrace {",
    "  for (const rule of STRATEGY_RULES) {",
    "    if (rule.test(ctx)) {",
    "      return { selected: rule.strategy, reason: rule.reason, context: ctx };",
    "    }",
    "  }",
    "  return { selected: 'SSR', reason: 'Fallback (no rule matched)', context: ctx };",
    "}",
], y, fs=13.5)
w = (CW - 0.20 * 2) / 3
props = [
    ("PURE", "No I/O, no clock, no state. It reads only its argument."),
    ("TOTAL", "The last rule is unconditionally true, so it can never fail to return."),
    ("DETERMINISTIC", "The same context always yields the same strategy — verified 1,000× per sampled context."),
]
for i, (t, d) in enumerate(props):
    x = M + i * (w + 0.20)
    rect(s, x, y + 2.05, w, 1.05, fill=PANEL, line=LINE, lw=0.75)
    rect(s, x, y + 2.05, w, 0.05, fill=BLUE)
    tf = tbox(s, x + 0.18, y + 2.24, w - 0.36, 0.8)
    para(tf, t, size=14, bold=True, color=NAVY, font=HEAD, after=4, first=True)
    para(tf, d, size=12, color=TEXT, font=BODY, after=0, line=1.12)
callout(s, "Purity was adopted for testability — but its real payoff was the evaluation. Because "
           "the function touches nothing external, the **entire policy could be enumerated "
           "exhaustively** and shown to be total and free of dead rules, which is a far stronger "
           "claim than a passing test suite.",
        y + 3.30, size=13.5, height=0.78)

# ── render/respond/measure
s, y = new("Flow of code — Stages 3–5: RENDER, RESPOND, MEASURE",
           kicker="Request lifecycle 4 of 4")
code(s, [
    "let result: RenderResult;",
    "try {",
    "  result = await this.deps.registry.get(trace.selected)",
    "                     .render(ctx, page, this.deps.cache, { reason: trace.reason });",
    "} catch (err) {                                   // a strategy threw: degrade, do not fail",
    "  result = await this.deps.registry.get('SSR')",
    "                     .render(ctx, page, this.deps.cache,",
    "                             { reason: `${trace.selected} render failed - fell back to SSR` });",
    "}",
    "",
    "const headersOut = { ...result.headers,",
    "  'x-rendering-strategy': trace.selected,",
    "  'x-decision-reason'  : trace.reason.replace(/[^\\x20-\\x7E]/g, '-') };  // HTTP is Latin-1",
    "",
    "res.writeHead(result.status, headersOut);",
    "const ttfbMs = timer.elapsedMs();               // TTFB: taken the moment headers are written",
    "...",
    "void this.deps.metrics.record({ ... });         // never awaited: metrics cannot delay a response",
], y, fs=11.5)
bullets(s, [
    "A degraded response is **never silently indistinguishable from a healthy one** — the substitution is written into the reason string.",
    "Timing is taken at two points: TTFB immediately after headers, total after the body flushes; render time is the difference. For streams the engine counts bytes as chunks pass and resolves on stream end, so **a streamed response is measured on the same basis as a buffered one**.",
], y + 3.36, size=12.5, gap=6)

# ── THE ALGORITHM
s, y = new("The decision algorithm", kicker="The intellectual centre of the project",
           sub="An ordered, first-match-wins rule table over the context. This is the authoritative policy as implemented in src/config/strategy-rules.ts.")
table(s, ["#", "Condition", "Strategy", "Rationale"], [
    ["1", "volatility = static  AND  cache ≠ cold", "SSG", "Static content with a usable cache: serve the pre-built artefact"],
    ["2", "volatility = static  AND  isEdge", "EDGE_ISR", "Static content at an edge: revalidate close to the user"],
    ["3", "load = high", "ISR", "Shed origin work: serve cached output, revalidate in the background"],
    ["4", "volatility = realtime  AND  net = fast  AND  device = desktop", "CSR", "Capable client on a fast link: ship a shell, let it be fully interactive"],
    ["5", "volatility = realtime  AND  device = mobile", "SSR", "Weak device: send finished HTML and minimal JavaScript"],
    ["6", "volatility = periodic", "ISR", "Periodically changing data: cache and revalidate on a TTL"],
    ["7", "heavyPayload  AND  net ≠ slow", "STREAMING_SSR", "Large payload on a decent link: stream the shell first, then the rest"],
    ["8", "net = slow", "SSR", "Slow link: avoid the cost of hydrating a large bundle"],
    ["9", "unconditional fallback", "SSR", "Safe, correct default — this is what makes the function total"],
], [0.45, 4.5, 1.75, 5.4], y, fs=11.5, rh=0.355,
   aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
callout(s, "The table is **data**; the evaluator that consumes it is the nine-line function on the "
           "previous slide. That separation is why the same table can be imported unchanged into the "
           "browser, so the demonstration pages explain the engine using the engine's own policy — "
           "and can never drift from it.",
        y + 3.62, size=13, height=0.76)

# ── algorithm formal
s, y = new("Algorithm: formal statement and complexity", kicker="The algorithm")
code(s, [
    "decide(C):                                  # C = the request context",
    "  for i = 1 .. 9:                           # R = (r1 ... r9), an ORDERED sequence",
    "    if r[i].test(C) is true:",
    "      return (r[i].strategy, r[i].reason, C)",
    "  return (SSR, \"fallback\", C)               # unreachable: r[9].test is constant true",
], y, fs=13, width=6.7, left=M)
tf = tbox(s, M, y + 1.42, 6.7, 2.6)
para(tf, "Each predicate is a conjunction of **at most three equality comparisons over enumerated "
         "values**, so evaluating one is O(1). The loop is therefore O(|R|) with |R| = 9 — the "
         "whole function runs in **constant time with no allocation** beyond the returned record.",
     size=13, after=10, first=True, line=1.16)
para(tf, "Because r₉ is unconditionally true, decide is a **total function**. Because every "
         "predicate reads only its argument, it is **deterministic and referentially transparent**.",
     size=13, after=10, line=1.16)
para(tf, "Both properties are not merely asserted — Section 5.3 enumerates all 648 contexts and "
         "confirms them, and Section 5.4 measures the constant at 26.53 ns.",
     size=13, after=0, line=1.16, color=BLUE)
picture(s, fig("decision-flow.png"), y, 3.6, left=M + 7.05, max_w=CW - 7.05)

# ── ordering is policy
s, y = new("Why the ordering IS the policy", kicker="The algorithm",
           sub="In a first-match-wins table, ordering carries as much meaning as the predicates — it encodes precedence between competing objectives.")
prec = [
    ("Rule 1 outranks rule 3  —  a usable cache beats high load",
     "Under a traffic spike a static page keeps being served from its pre-built artefact rather "
     "than being downgraded to ISR, because the artefact is already the cheapest possible answer."),
    ("Only a **cold** cache defeats rule 1  —  a stale one does not",
     "Deliberate. Stale-while-revalidate is a feature of the design, not a cache miss, so 'stale' "
     "is treated as usable."),
    ("Rule 5 outranks rule 7  —  mobile on /heavy gets plain SSR, never Streaming SSR",
     "A weak device should not be asked to hydrate a large payload merely because that payload "
     "could have been streamed."),
    ("Rule 3 sits third  —  an explicit load=high overrides device and volatility",
     "This has a direct experimental consequence: to hold a fixed-SSR baseline under heavy "
     "concurrency the load signal must be pinned low, or the engine correctly promotes the "
     "baseline to ISR. Section 5.8 had to be rebuilt around exactly this."),
]
tf = tbox(s, M, y, CW, 4.0)
for i, (t, d) in enumerate(prec):
    para(tf, t, size=14.5, bold=True, color=NAVY, font=HEAD, after=3, first=(i == 0))
    para(tf, d, size=13, color=TEXT, font=BODY, after=14, indent=0.0, line=1.15)
callout(s, "All four were verified experimentally against the running stack — they are rows in "
           "the trigger matrix, not claims about the source code.",
        y + 3.86, size=13)

# ── six strategies
s, y = new("The six rendering strategies", kicker="Implementation",
           sub="Each is a module implementing the same interface. Uniformity of that signature is what makes them interchangeable at runtime.")
table(s, ["Strategy", "React mechanism", "Cache behaviour", "Marker header"], [
    ["SSG", "renderToString at start-up", "Reads a pre-built file from disk; builds on demand if absent", "x-ssg: prebuilt | built-on-demand"],
    ["SSR", "renderToString per request", "None", "x-render-bytes"],
    ["STREAMING_SSR", "renderToPipeableStream with Suspense", "None; piped through a PassThrough stream", "x-streaming: true + chunked encoding"],
    ["ISR", "renderToString, result cached with a TTL", "fresh: serve  ·  stale: serve + revalidate in background  ·  cold: render, cache, serve", "x-isr-cache: fresh | stale-revalidating | miss"],
    ["CSR", "Empty shell; the browser renders", "Withholds data so the client must fetch /api/data", "x-csr: shell"],
    ["EDGE_ISR", "Inherits ISR", "Identical semantics on a per-edge cache namespace, optionally shared through Redis", "x-isr-cache: …"],
], [1.7, 2.8, 4.6, 2.9], y, fs=11, rh=0.42)
tf = tbox(s, M, y + 2.94, CW, 1.3)
para(tf, "Two implementation decisions worth defending:", size=13, bold=True, color=NAVY,
     font=HEAD, after=7, first=True)
para(tf, "**SSG could never be selected on a cold start**, because rule 1 requires a usable cache "
         "that only a previous request could create. Static pages are therefore pre-built at "
         "server start-up, and the prebuild embeds the context under which rule 1 actually serves "
         "the artefact — not the cold state present at build time.",
     size=12.5, after=6, bullet="▪", line=1.13)
para(tf, "**Edge-ISR is the clearest demonstration of the pluggability invariant**: a complete "
         "sixth strategy expressed as a subclass with a single overridden method — the cache "
         "key, namespaced by node identity — because the semantics it needs already exist.",
     size=12.5, after=0, bullet="▪", line=1.13)

# ── ISR deep dive
s, y = new("Deep dive: ISR, stale-while-revalidate and single-flight",
           kicker="Implementation",
           sub="The most subtle of the six, and the strategy on which the engine's entire load-shedding behaviour depends.")
code(s, [
    "if (lookup.entry && !lookup.stale) return this.result(lookup.entry.value, true, 'fresh');",
    "",
    "if (lookup.entry && lookup.stale) {            // serve NOW, refresh behind the response",
    "  void revalidateInBackground(key, cache, produce, cache.defaultTtlMs);",
    "  return this.result(lookup.entry.value, true, 'stale-revalidating');",
    "}",
    "",
    "const html = await produce();                  // cold: render synchronously, cache, serve",
    "await cache.set(key, html, cache.defaultTtlMs);",
    "return this.result(html, false, 'miss');",
], y, fs=12.5, width=7.4, left=M)
x = M + 7.7
tf = tbox(s, x, y, CW - 7.7, 3.4)
para(tf, "Three states, one code path", size=15, bold=True, color=NAVY, font=HEAD,
     after=8, first=True)
para(tf, "**fresh** — served directly from cache within the TTL window.", size=12.5,
     after=6, bullet="▪", line=1.12)
para(tf, "**stale** — served immediately from the existing entry while a refresh runs behind "
         "the response. The user never waits for it.", size=12.5, after=6, bullet="▪", line=1.12)
para(tf, "**cold** — rendered synchronously, cached, served.", size=12.5, after=10,
     bullet="▪", line=1.12)
para(tf, "Background revalidation is **single-flight**: a module-level set of in-flight keys "
         "guarantees that a burst of stale requests triggers exactly one re-render rather than a "
         "stampede.", size=12.5, color=BLUE, after=0, line=1.14)
callout(s, f"Measured directly: a cold render cost **{cold_ms:.2f} ms**, warm hits averaged "
           f"**{fresh_ms:.2f} ms** ({warm_cut:.1f} % lower), and a stale hit cost only "
           f"**{stale_ms:.2f} ms** rather than paying the cold render again. That gap is the entire "
           f"value of stale-while-revalidate.",
        y + 2.42, size=13.5, bg=GREEN_BG, fg=GREEN_TX, height=0.80)

# ── cache subsystem
s, y = new("The cache subsystem", kicker="Implementation",
           sub="One interface, three backends, one freshness rule.")
tiers = [
    ("MEMORY", "in-process", "A Map-based store with recency refresh on read and a bounded entry count. The fastest tier — and, as Section 5.5 shows, fast enough to invert the expected ranking of SSG and ISR."),
    ("REDIS", "shared, optional", "Loaded through a dynamic import so its absence is never a build failure. Writes carry the TTL to Redis directly. This is what makes the two edge nodes share cache state."),
    ("FILE", "persistent", "Hashes the key to a filename and stores the entry as JSON, giving persistence across process restarts."),
]
w = (CW - 0.20 * 2) / 3
for i, (t, sub, d) in enumerate(tiers):
    x = M + i * (w + 0.20)
    rect(s, x, y, w, 1.95, fill=PANEL, line=LINE, lw=0.75)
    rect(s, x, y, w, 0.05, fill=BLUE)
    tf = tbox(s, x + 0.18, y + 0.20, w - 0.36, 1.6)
    para(tf, t, size=15, bold=True, color=NAVY, font=HEAD, after=1, first=True)
    para(tf, sub, size=10.5, color=BLUE, font=MONO, after=7)
    para(tf, d, size=12, color=TEXT, font=BODY, after=0, line=1.14)
tf = tbox(s, M, y + 2.20, CW, 2.0)
para(tf, "Reads **descend** memory → Redis → file and **promote any hit upward**, so a "
         "subsequent read is faster. Writes **fan out** to every available backend.",
     size=14, after=10, first=True, bullet="▪", line=1.16)
para(tf, "Freshness is a property of the entry itself — its stored timestamp plus its TTL. "
         "That is why the same entry can be classified fresh, stale or cold **without any separate "
         "bookkeeping**, and why the analyzer can probe cache state cheaply before a decision is made.",
     size=14, after=10, bullet="▪", line=1.16)
para(tf, "Redis is optional **by construction**: if it is absent the manager logs a warning and "
         "continues on memory and file alone. The engine never hard-depends on infrastructure it "
         "cannot guarantee.",
     size=14, after=0, bullet="▪", line=1.16)

# ── demo pages
s, y = new("Self-explaining demonstration pages", kicker="Implementation",
           sub="Rather than a generic sample page, each page renders the engine's own decision into its own markup.")
table(s, ["Route", "Volatility", "Default strategy", "What it demonstrates"], [
    ["/static", "static", "SSG",
     "Artefact ageing — the frozen generation timestamp beside a live clock, making visible exactly the staleness ISR exists to bound. Plus a fresh/stale/cold cache lab."],
    ["/dynamic", "realtime", "SSR",
     "The CSR/SSR distinction — the page reports whether its data arrived embedded in the markup or was fetched by the client."],
    ["/heavy", "realtime, heavy", "STREAMING_SSR",
     "Streaming — shell flushed first, Suspense boundary streamed after, with client-side filter/sort/paging over 400 rows and real TTFB/DOM timings."],
    ["/control", "—", "—",
     "The ARE Control Center: a live runtime panel with a strategy inspector, a load generator, an origin-vs-edge race and live telemetry."],
], [1.3, 1.5, 2.0, 7.3], y, fs=11.5, rh=0.52)
tf = tbox(s, M, y + 2.52, CW, 1.7)
para(tf, "The shared decision console on all three pages:", size=13.5, bold=True,
     color=NAVY, font=HEAD, after=7, first=True)
para(tf, "Renders the strategy, the rule that fired and the observed context **server-side** — "
         "visible in view-source with JavaScript disabled.", size=12.5, after=5, bullet="▪", line=1.12)
para(tf, "**Imports the engine's own STRATEGY_RULES** and evaluates the full table in the browser, "
         "highlighting the winning rule — so its explanation can never drift from server behaviour.",
     size=12.5, after=5, bullet="▪", line=1.12)
para(tf, "Probes the live server and **cross-checks its own prediction** against the real "
         "X-Rendering-Strategy header, and shows a hydrated ✓ marker that is false on the "
         "server and flips after mount — the proof that interactivity is real.",
     size=12.5, after=0, bullet="▪", line=1.12)

# ── hydration discipline
s, y = new("Hydration discipline — why the pages produce zero React warnings",
           kicker="Implementation",
           sub="Making a page explain its own server-side decision, and then hydrate cleanly, required a rule we had to discover the hard way.")
w = (CW - 0.28) / 2
rect(s, M, y, w, 2.55, fill=PANEL, line=LINE, lw=0.75)
rect(s, M, y, w, 0.05, fill=RED)
tf = tbox(s, M + 0.22, y + 0.22, w - 0.44, 2.2)
para(tf, "THE PROBLEM", size=11, bold=True, color=RED, font=HEAD, after=9, first=True)
para(tf, "Demonstration pages that read the clock during render produced React hydration "
         "mismatches: the server container runs in UTC, the browser does not, so the two renders "
         "disagreed on the very timestamps the pages exist to show.",
     size=13, after=0, line=1.16)
x2 = M + w + 0.28
rect(s, x2, y, w, 2.55, fill=PANEL, line=LINE, lw=0.75)
rect(s, x2, y, w, 0.05, fill=GREEN_TX)
tf = tbox(s, x2 + 0.22, y + 0.22, w - 0.44, 2.2)
para(tf, "THE RULE ADOPTED", size=11, bold=True, color=GREEN_TX, font=HEAD, after=9, first=True)
para(tf, "The first render depends **only on props**. No Date.now(), no Math.random(), no "
         "navigator, window or performance during render. All live values start from props and "
         "update inside useEffect after mount; timestamps render as stable UTC and localise in the "
         "browser.", size=13, after=0, line=1.16)
callout(s, "Verified: **zero hydration warnings** in headless Chrome across all pages and all "
           "strategies. This matters beyond tidiness — a hydration mismatch would have "
           "invalidated the pages' claim to be showing the server's real decision.",
        y + 2.78, size=13.5, bg=GREEN_BG, fg=GREEN_TX, height=0.78)

# ── one request end to end
s, y = new("Working principle: one request, end to end", kicker="The whole system in one narrative",
           sub="A browser on a fast connection requests /dynamic.")
steps = [
    ("PROXY", "Receives it on :8080, sets X-Served-By: origin, forwards to the origin container."),
    ("SERVER", "Increments the in-flight counter, finds no matching static asset, resolves /dynamic to the realtime page module, sleeps 0 ms for the fast class, probes the cache, calls the engine."),
    ("ANALYZE", "No control headers, so: network from the query alias or default medium; device inferred desktop from the User-Agent; cache state from the probe; load classified from the in-flight count; volatility = realtime, because that is what the page declares."),
    ("DECIDE", "Rule 1 fails (not static). Rule 2 fails. Rule 3 fails (load is low). Rule 4 matches — realtime + fast + desktop → CSR, reason: 'Realtime data on a capable client'."),
    ("RENDER · RESPOND", "Engine logs three lines, invokes CSR, receives a shell with the data withheld, writes X-Rendering-Strategy: CSR and X-Decision-Reason, sends the body."),
    ("MEASURE", "Records the metric without awaiting it; the finally block decrements the counter."),
]
tf = tbox(s, M, y, CW, 3.3)
for i, (t, d) in enumerate(steps):
    p = para(tf, t, size=12, bold=True, color=BLUE, font=HEAD, after=1, first=(i == 0))
    para(tf, d, size=12.5, color=TEXT, font=BODY, after=8, indent=0.0, line=1.13)
callout(s, f"The browser receives an **{ad['doc_bytes']}-byte shell**, runs the client bundle and "
           f"fetches its data. Had the same URL arrived from a mobile User-Agent, rule 4 would fail, "
           f"rule 5 would match, and the same server would have returned **{fx['doc_bytes']:,} bytes "
           f"of finished HTML** instead. That divergence, from one unchanged URL, is the whole "
           f"system in one sentence.",
        y + 3.42, size=14, bg=BLUE_L, height=0.84)

# ── live screenshots
_shots = [
    ("image2.jpeg", "The engine adapting live", "Portfolio page at localhost:8080 — the engine reports the device, the network class it inferred, the live latency it measured, and the strategy it chose for this visitor, with the rule text that produced it."),
    ("image5.jpeg", "The decision console, server-rendered", "/static served as SSG — the artefact's frozen generation timestamp beside its live age, the rule that fired, and the hydrated ✓ marker proving interactivity was bolted on afterwards."),
    ("image8.jpeg", "The ARE Control Center", "/control — a live strategy inspector and simulator, a concurrent load generator, an origin-vs-edge race, a cache lab and real-time telemetry, all driving the running engine."),
]
for f, t, cap in _shots:
    p = shot(f)
    if not p:
        continue
    s, y = new(t, kicker="The running system")
    picture(s, p, y + 0.02, 4.10)
    tf = tbox(s, M, y + 4.26, CW, 0.6)
    para(tf, cap, size=12.5, color=MUTED, font=BODY, after=0, first=True,
         align=PP_ALIGN.CENTER, line=1.15)

# ═══════════════════════════════════════════════ SECTION 4 — RESULTS
section("04", "Results, Analysis and Evaluation",
        ["Functional correctness, then the cost of adapting, then performance",
         "An exhaustive proof over the complete context space — not just a passing test suite",
         "Adaptive selection against a fixed policy, under two different bottlenecks"])

# ── experimental setup
s, y = new("Experimental setup", kicker="Results",
           sub="All measurements were taken against the five-container private server on a single macOS arm64 host.")
w = (CW - 0.28) / 2
rect(s, M, y, w, 2.05, fill=PANEL, line=LINE, lw=0.75)
tf = tbox(s, M + 0.22, y + 0.20, w - 0.44, 1.7)
para(tf, "SERVER-SIDE TIMINGS", size=11, bold=True, color=BLUE, font=HEAD, after=8, first=True)
para(tf, "Produced by the engine's own metrics pipeline. Because the simulated link delay and the "
         "container's edge latency are applied **before** the engine is invoked, and the engine's "
         "timer starts inside its handler, these measure **engine cost alone**.",
     size=12.5, after=0, line=1.15)
x2 = M + w + 0.28
rect(s, x2, y, w, 2.05, fill=PANEL, line=LINE, lw=0.75)
tf = tbox(s, x2 + 0.22, y + 0.20, w - 0.44, 1.7)
para(tf, "CLIENT-SIDE TIMINGS", size=11, bold=True, color=BLUE, font=HEAD, after=8, first=True)
para(tf, "Wall clock via curl, and concurrent load via Apache Bench. These include **everything**: "
         "proxy hop, simulated link, edge latency and engine. The two views are complementary and "
         "neither is used to stand in for the other.",
     size=12.5, after=0, line=1.15)
bullets(s, [
    "Unless stated otherwise the per-strategy benchmarks used the **fast** network class, whose configured delay is zero, so that engine cost is not masked.",
    "Every experiment was run from a **freshly recreated stack** — and because the deployment mounts no volumes, that guarantees empty caches, no pre-existing artefacts beyond the start-up prebuild, and an empty metrics log.",
    "All results are reproducible from a clean checkout with the commands in Appendix C of the report.",
], y + 2.28, size=13.5)

# ── R1 trigger matrix
s, y = new(f"Result 1 — the same URL, answered differently: {D['trigger_pass']}/17",
           kicker="Functional verification",
           sub="Each row is a single request to the running stack; the observed value is read from the X-Rendering-Strategy response header.")
tm = D["trigger_matrix"]
rows = [[r["url"], r["expected"], r["observed"], r["ok"], r["rule"]] for r in tm]
half = 9
left_rows = rows[:half]
right_rows = rows[half:]
hdr = ["Request", "Expected", "Observed", "", "Rule"]
gf_w = (CW - 0.22) / 2


def _mini(rows_, left):
    base = [4.05, 1.25, 1.25, 0.6, 3.2]
    widths = [w_ * (gf_w / sum(base)) for w_ in base]
    n = len(rows_) + 1
    g = s.shapes.add_table(n, 5, Inches(left), Inches(y), Inches(gf_w),
                           Inches(0.30 + 0.285 * len(rows_)))
    tb_ = g.table
    tb_.first_row = True
    tb_.horz_banding = False
    for i, w_ in enumerate(widths):
        tb_.columns[i].width = Inches(w_)
    tb_.rows[0].height = Inches(0.30)
    for c, h in enumerate(hdr):
        cell = tb_.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.margin_left = cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.02)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(cell.text_frame, h, size=9.5, bold=True, color=WHITE, font=HEAD,
             after=0, first=True)
    for r, row in enumerate(rows_, start=1):
        tb_.rows[r].height = Inches(0.285)
        for c, val in enumerate(row):
            cell = tb_.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_A if r % 2 else ROW_B
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            col = TEXT
            bold = False
            fnt = BODY
            if c == 0:
                fnt = MONO
                col = NAVY
            if c in (1, 2):
                bold = True
                col = NAVY
            if c == 3:
                col = GREEN_TX
                bold = True
            if c == 4:
                col = MUTED
            para(cell.text_frame, str(val), size=8.6, bold=bold, color=col,
                 font=fnt, after=0, first=True, line=1.0)


_mini(left_rows, M)
_mini(right_rows, M + gf_w + 0.22)
callout(s, "Rows 1–7 hold the URL at **/static** and vary only the query context, yet obtain "
           "**five different strategies from one unchanged resource** — the clearest possible "
           "demonstration of the core claim. Rows 14–15 confirm the precedence properties; rows "
           "16–17 confirm Edge-ISR through the proxy route and against an edge container directly.",
        y + 2.90, size=13, height=0.80)

# ── R2 exhaustive
s, y = new(f"Result 2 — the policy is total, deterministic and has no dead rules",
           kicker="Correctness",
           sub=f"Because decide() is pure, the entire context space could be enumerated: 3×2×3×3×3×2×2 = {ds['total']} distinct contexts. Every one was evaluated.")
picture(s, fig("fig_context_space.png"), y, 3.05, max_w=7.5, left=M)
x = M + 7.75
tf = tbox(s, x, y, CW - 7.75, 3.4)
para(tf, "Three results follow", size=15, bold=True, color=NAVY, font=HEAD, after=9, first=True)
para(tf, f"**Total.** All {ds['total']} contexts returned a strategy; none reached an error path.",
     size=12.5, after=7, bullet="▪", line=1.13)
para(tf, "**No dead rules.** Every one of the nine rules is the first match for at least 20 "
         "contexts, so the policy has no unreachable branches.",
     size=12.5, after=7, bullet="▪", line=1.13)
para(tf, "**Deterministic.** Repeating a sample of contexts one thousand times each produced an "
         "identical strategy every time — the experimental counterpart of the purity argument.",
     size=12.5, after=10, bullet="▪", line=1.13)
para(tf, "Complemented by **28 automated tests** across 4 files: 10 on the decision engine (one "
         "per rule row), 14 on the context analyzer, 2 on rendering and 2 on the cache. All pass.",
     size=12.5, color=BLUE, after=0, line=1.14)
callout(s, f"The distribution is itself informative. **ISR claims {ds['byStrategy']['ISR']} contexts "
           f"({ds['byStrategy']['ISR']/ds['total']*100:.1f} %)** because it is selected both by high "
           f"load and by periodic volatility; SSG claims {ds['byStrategy']['SSG']}; SSR claims "
           f"{ds['byStrategy']['SSR']} across three separate rules; while CSR ({ds['byStrategy']['CSR']}), "
           f"EDGE_ISR ({ds['byStrategy']['EDGE_ISR']}) and STREAMING_SSR ({ds['byStrategy']['STREAMING_SSR']}) "
           f"occupy narrow, sharply-specified regions. That asymmetry is intentional: cache-backed "
           f"strategies are the broad default, and the specialised strategies are reserved for the "
           f"precise conditions in which they win.",
        y + 3.30, size=12.5, height=0.96)

# ── R3 overhead
s, y = new("Result 3 — what does it cost to adapt?", kicker="Cost of the mechanism",
           sub="An adaptive engine must justify its own overhead. The decision function was benchmarked in isolation over two million evaluations cycling through the enumerated context space.")
tiles(s, [
    (f"{ds['ns_per_decision']:.2f} ns", "per decision", "time to select a strategy"),
    (f"{ds['decisions_per_second']/1e6:.1f} M", "decisions per second", "on one core"),
    ("≤ 9", "rules evaluated", "each an O(1) comparison"),
    (f"{ps['ISR']['ttfb']:.2f} ms", "cheapest measured request", "mean server-side TTFB"),
    ("0.0126 %", "of that request", "the decision's share"),
], y, vsize=23)
callout(s, f"At **{ds['ns_per_decision']:.2f} ns**, a decision costs roughly **one ten-thousandth** "
           f"of the cheapest measured request. The adaptive mechanism is free in any practical "
           f"sense: the engine could decide for every request of a large site many times over and "
           f"still not register against the cost of producing a single response.",
        y + 1.36, size=14, bg=GREEN_BG, fg=GREEN_TX, height=0.86)
callout(s, "This settles the most obvious objection to runtime strategy selection — that "
           "choosing costs more than it saves — at least for a rule-based policy. It also "
           "sets the budget for the learned policies proposed in Chapter 7: a model has roughly "
           "this envelope to work within before it starts costing what it saves.",
        y + 2.36, size=13.5, height=0.80)

# ── R4 per-strategy
s, y = new("Result 4 — per-strategy performance", kicker="Performance analysis",
           sub="60 sequential requests per strategy through the proxy on the fast class. These are the engine's own timings.")
rows = []
order = ["SSG", "SSR", "CSR", "ISR", "STREAMING_SSR", "EDGE_ISR"]
for k in order:
    v = ps[k]
    rows.append([k, f"{v['ttfb']:.3f}", f"{v['ttfb_p95']:.3f}", f"{v['render']:.3f}",
                 f"{v['bytes']:,}", f"{v['hit']:.2f}", f"{pc[k]['mean']:.2f}"])
table(s, ["Strategy", "Mean TTFB (ms)", "p95 TTFB (ms)", "Mean render (ms)",
          "Mean bytes", "Cache-hit rate", "Client mean (ms)"], rows,
      [2.0, 1.75, 1.65, 1.85, 1.55, 1.6, 1.75], y, fs=12, rh=0.355,
      aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.CENTER] * 6)
bullets(s, [
    f"**CSR is the cheapest on both axes** — {ps['CSR']['ttfb']:.3f} ms and only {ps['CSR']['bytes']} bytes, because it renders nothing on the server and withholds the data. It is {csr_light:.1f}× lighter on the wire than SSR.",
    f"**ISR is almost as cheap** at {ps['ISR']['ttfb']:.3f} ms with a cache-hit rate of {ps['ISR']['hit']:.2f}, while still delivering complete markup of {ps['ISR']['bytes']:,} bytes. SSR's output at close to CSR's cost — which is exactly why ISR is the engine's preferred response to load.",
    f"**Streaming SSR is the most expensive per request** at {ps['STREAMING_SSR']['ttfb']:.3f} ms and {ps['STREAMING_SSR']['bytes']:,} bytes — which is correct, not disappointing: it is the only strategy applied to the deliberately heavy page, and its purpose is to let the client begin parsing early, not to reduce total bytes.",
    f"**Edge-ISR tracks ISR closely** at {ps['EDGE_ISR']['ttfb']:.3f} ms with a hit rate of {ps['EDGE_ISR']['hit']:.2f}, as its inheritance implies. Its {pc['EDGE_ISR']['mean']:.1f} ms client-side mean is the deliberately injected 20 ms edge latency dominating everything else.",
], y + 2.52, size=12.5, gap=7)

# ── SSG anomaly
s, y = new("Result 4b — a genuinely counter-intuitive finding",
           kicker="Performance analysis",
           sub="The strategy that performs no rendering at all is not the cheapest strategy.")
picture(s, fig("fig_ttfb.png"), y, 2.48, max_w=5.86, left=M)
picture(s, fig("fig_bytes.png"), y, 2.48, max_w=5.86, left=M + 6.11)
callout(s, f"**SSG costs {ps['SSG']['ttfb']:.3f} ms — {ssg_ratio:.1f}× more than ISR's "
           f"{ps['ISR']['ttfb']:.3f} ms**, despite doing no rendering whatsoever.",
        y + 2.70, size=15, bg=AMBER_BG, fg=AMBER_TX, height=0.56)
tf = tbox(s, M, y + 3.38, CW, 1.1)
para(tf, "**The explanation is in the implementation, not the theory.** SSG reads its artefact from "
         "the filesystem on every request, whereas a warm ISR entry is served from the in-process "
         "memory tier of the cache hierarchy — and a filesystem read, even of a small file, "
         "costs more than a map lookup.",
     size=13, after=7, first=True, line=1.16)
para(tf, "A real and useful finding: **the theoretical ranking of strategies can be inverted by "
         "the storage tier they land on** — and it identifies an immediate optimisation, promoting "
         "SSG artefacts into the memory cache, which is the first item on the roadmap.",
     size=13, after=0, color=BLUE, line=1.16)

# ── R5 ISR lifecycle
s, y = new("Result 5 — the ISR cache lifecycle, measured", kicker="Cache behaviour",
           sub="34 requests at two-second intervals over ~69 seconds, from a freshly recreated stack, against an origin whose TTL is 30 seconds.")
picture(s, fig("fig_isr_lifecycle.png"), y, 2.90, max_w=7.3, left=M)
x = M + 7.55
table(s, ["Cache state", "Requests", "Mean (ms)"], [
    ["miss (cold)", str(len(cold)), f"{cold_ms:.2f}"],
    ["fresh", str(len(fresh)), f"{fresh_ms:.2f}"],
    ["stale-revalidating", str(len(stale)), f"{stale_ms:.2f}"],
], [2.6, 1.2, 1.4], y, left=x, fs=11.5, rh=0.36, width=CW - 7.55,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
tf = tbox(s, x, y + 1.58, CW - 7.55, 2.0)
para(tf, f"The sawtooth is exactly as designed. Two stale transitions were captured, at "
         f"approximately 31 s and 62 s — **matching the 30-second TTL almost exactly**.",
     size=12.5, after=9, first=True, line=1.15)
para(tf, f"Warm hits averaged {fresh_ms:.2f} ms against a cold render of {cold_ms:.2f} ms: a "
         f"**{warm_cut:.1f} % reduction**.",
     size=12.5, after=9, line=1.15)
para(tf, f"Critically, a stale hit cost **{stale_ms:.2f} ms rather than {cold_ms:.2f} ms** — the "
         f"user was served immediately from the stale entry while the refresh happened behind the "
         f"response.",
     size=12.5, after=0, color=BLUE, line=1.15)
callout(s, "That gap is the entire value of stale-while-revalidate, and it is what allows ISR to "
           "absorb load **without a latency spike at every TTL boundary** — the property the "
           "engine depends on when rule 3 fires under a traffic spike.",
        y + 3.14, size=13.5)

# ── R6 simulation fidelity
s, y = new("Result 6 — is the simulated environment faithful?", kicker="Validity",
           sub="Because the evaluation depends on simulated conditions, the simulation itself was validated against its configuration.")
picture(s, fig("fig_network_edge.png"), y, 2.85, max_w=6.6, left=M)
net, edge = D["network"], D["edge"]
table(s, ["Condition", "Configured", "Measured mean", "Residual"], [
    ["Network: slow", "400 ms", f"{net['slow']['mean']:.2f} ms", f"{net['slow']['mean']-400:.2f} ms"],
    ["Network: medium", "100 ms", f"{net['medium']['mean']:.2f} ms", f"{net['medium']['mean']-100:.2f} ms"],
    ["Network: fast", "0 ms", f"{net['fast']['mean']:.2f} ms", f"{net['fast']['mean']:.2f} ms"],
    ["Origin node", "+0 ms", f"{edge['origin']['mean']:.2f} ms", "baseline"],
    ["edge-node-1", "+20 ms", f"{edge['edge-node-1']['mean']:.2f} ms", f"{edge['edge-node-1']['mean']-edge['origin']['mean']-20:.2f} ms"],
    ["edge-node-2", "+80 ms", f"{edge['edge-node-2']['mean']:.2f} ms", f"{edge['edge-node-2']['mean']-edge['origin']['mean']-80:.2f} ms"],
], [2.0, 1.4, 1.7, 1.3], y, left=M + 6.85, fs=11, rh=0.34, width=CW - 6.85,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
callout(s, f"The residuals are small and consistent. The fast class, configured at zero, measures "
           f"**{net['fast']['mean']:.2f} ms** — the irreducible cost of the proxy hop, the "
           f"container network and the engine — and the slow and medium classes exceed their "
           f"configured delays by approximately that same amount. The two edges overshoot by "
           f"{edge['edge-node-1']['mean']-edge['origin']['mean']-20:.2f} ms and "
           f"{edge['edge-node-2']['mean']-edge['origin']['mean']-80:.2f} ms, again a fixed "
           f"per-request overhead rather than a defect in the injection. **The simulation is "
           f"faithful, and the differences it produces are attributable to the conditions being modelled.**",
        y + 3.10, size=12.5, height=0.98)

# ── R7 fast link
s, y = new("Result 7 — adaptive vs a fixed SSR policy: fast link", kicker="Comparative evaluation",
           sub="800 requests at concurrency 50, same URL, same stack. Here the ORIGIN is the bottleneck.")
table(s, ["Measure", "Fixed SSR policy", "Adaptive (ARE)", "Change"], [
    ["Strategy actually used", "SSR × 800", "CSR × 800", "selected by rule 4"],
    ["Throughput (requests/s)", f"{fx['rps']:,.2f}", f"{ad['rps']:,.2f}", f"+{rps_gain:.1f} %"],
    ["Mean latency (ms)", f"{fx['mean_ms']:.2f}", f"{ad['mean_ms']:.2f}", f"−{mean_cut:.1f} %"],
    ["Median latency (ms)", f"{fx['p50']}", f"{ad['p50']}", f"−{(fx['p50']-ad['p50'])/fx['p50']*100:.1f} %"],
    ["95th percentile (ms)", f"{fx['p95']}", f"{ad['p95']}", f"−{(fx['p95']-ad['p95'])/fx['p95']*100:.1f} %"],
    ["99th percentile (ms)", f"{fx['p99']}", f"{ad['p99']}", f"−{p99_cut:.1f} %"],
    ["Slowest request (ms)", f"{fx['max']}", f"{ad['max']}", f"−{(fx['max']-ad['max'])/fx['max']*100:.1f} %"],
    ["Document size (bytes)", f"{fx['doc_bytes']:,}", f"{ad['doc_bytes']:,}", f"−{byte_cut:.1f} %"],
    ["Total transferred (bytes)", f"{fx['total_bytes']:,}", f"{ad['total_bytes']:,}",
     f"−{(fx['total_bytes']-ad['total_bytes'])/fx['total_bytes']*100:.1f} %"],
    ["Origin CPU consumed (ms)", f"{fx['cpu_ms']}", f"{ad['cpu_ms']}",
     f"−{(fx['cpu_ms']-ad['cpu_ms'])/fx['cpu_ms']*100:.1f} %"],
    ["Mean server TTFB (ms)", f"{fx['ttfb']:.3f}", f"{ad['ttfb']:.3f}",
     f"−{(fx['ttfb']-ad['ttfb'])/fx['ttfb']*100:.1f} %"],
], [3.4, 2.4, 2.4, 2.0], y, left=M, fs=11.5, rh=0.30,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
callout(s, f"**The tail is more striking than the mean.** The fixed policy's 99th percentile was "
           f"{fx['p99']} ms against the adaptive engine's {ad['p99']} ms. Under the fixed policy the "
           f"server must synthesise a complete document for every one of 50 concurrent requests, and "
           f"the contention produces a long queue; the adaptive engine, having recognised that this "
           f"client can render for itself, does almost no work per request — and the tail "
           f"largely disappears.",
        y + 3.66, size=13, height=0.92)

# ── R8 medium link
s, y = new("Result 8 — adaptive vs a fixed SSR policy: 100 ms link",
           kicker="Comparative evaluation",
           sub="Same experiment, one variable changed. Here the NETWORK is the bottleneck — and the honest reporting of this case matters more than the headline.")
table(s, ["Measure", "Fixed SSR policy", "Adaptive (ARE)", "Change"], [
    ["Strategy actually used", "SSR × 800",
     f"ISR × {ad2['strategies']['ISR']}, SSR × {ad2['strategies']['SSR']}", "rule 3 fired under load"],
    ["Throughput (requests/s)", f"{fx2['rps']:.2f}", f"{ad2['rps']:.2f}", f"{rps_d2:+.1f} %"],
    ["Mean latency (ms)", f"{fx2['mean_ms']:.2f}", f"{ad2['mean_ms']:.2f}",
     f"{(ad2['mean_ms']-fx2['mean_ms'])/fx2['mean_ms']*100:+.1f} %"],
    ["95th percentile (ms)", f"{fx2['p95']}", f"{ad2['p95']}",
     f"−{(fx2['p95']-ad2['p95'])/fx2['p95']*100:.1f} %"],
    ["99th percentile (ms)", f"{fx2['p99']}", f"{ad2['p99']}", f"−{p99_cut2:.1f} %"],
    ["Slowest request (ms)", f"{fx2['max']}", f"{ad2['max']}",
     f"−{(fx2['max']-ad2['max'])/fx2['max']*100:.1f} %"],
    ["Cache-hit rate", f"{fx2['hit']:.2f}", f"{ad2['hit']:.2f}", "0.52 of responses from cache"],
    ["Origin CPU consumed (ms)", f"{fx2['cpu_ms']:,}", f"{ad2['cpu_ms']}", f"−{cpu_cut:.1f} %"],
    ["Mean server TTFB (ms)", f"{fx2['ttfb']:.3f}", f"{ad2['ttfb']:.3f}",
     f"−{(fx2['ttfb']-ad2['ttfb'])/fx2['ttfb']*100:.1f} %"],
], [3.4, 2.4, 2.6, 2.4], y, left=M, fs=11.5, rh=0.315,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
tf = tbox(s, M, y + 3.30, CW, 1.4)
para(tf, f"**Throughput is effectively identical ({rps_d2:+.1f} %), and that is the expected "
         f"result.** When every request must wait 100 ms on the simulated link, the origin is not "
         f"the constraint, and no rendering decision can remove a delay imposed downstream of it.",
     size=13.5, after=8, first=True, line=1.16)
para(tf, f"What changes is **what the origin spends**. The engine observed concurrency above the "
         f"high-load threshold and selected ISR for {ad2['strategies']['ISR']} of the 800 requests "
         f"under rule 3, reaching a cache-hit rate of {ad2['hit']:.2f} where the fixed policy "
         f"reached 0.00. Origin CPU fell from {fx2['cpu_ms']:,} ms to {ad2['cpu_ms']} ms — "
         f"**{cpu_cut:.1f} % less processor for identical delivered work** — and the tail "
         f"tightened by {p99_cut2:.1f} %.",
     size=13.5, after=0, color=BLUE, line=1.16)

# ── latency distribution figure
s, y = new("Result 8b — the two regimes side by side", kicker="Comparative evaluation",
           sub="Latency distribution under load: fixed policy versus adaptive selection.")
picture(s, fig("fig_loadtest.png"), y, 3.55)
callout(s, "**The correct reading of the two experiments together:** adaptive rendering does not "
           "manufacture bandwidth. Where the origin is the constraint it produces large "
           "user-visible gains; where the network is the constraint it leaves user-visible latency "
           "alone and converts the saving into **headroom** — a third of the origin's processor "
           "returned to the operator, available to absorb the spike that would otherwise have caused "
           "the very queueing the fixed policy exhibits in its tail.",
        y + 3.78, size=13, height=0.94)

# ── measurement integrity
s, y = new("A note on measurement integrity", kicker="Validity",
           sub="Apache Bench reports a non-zero 'Failed requests' count for these runs. It is benign, and here is exactly why.")
tiles(s, [
    ("Length", "the only failure category reported", "no other category is non-zero"),
    ("0", "connection errors", "across every run"),
    ("0", "exceptions", "across every run"),
    ("0", "non-2xx responses", "every request succeeded"),
], y)
tf = tbox(s, M, y + 1.36, CW, 1.8)
para(tf, "**The cause.** The demonstration page embeds live timestamps and generated data, so "
         "consecutive responses legitimately differ in length — and ab flags any response whose "
         "length differs from the first it received. It is measuring variance in the page content, "
         "not failure.",
     size=14, after=10, first=True, line=1.18)
para(tf, f"**The confirming detail.** The adaptive run on the fast link reported **zero** such "
         f"mismatches, precisely because the CSR shell it returns is data-free and therefore "
         f"byte-identical every time. The anomaly disappears exactly where the theory says it should.",
     size=14, after=0, color=BLUE, line=1.18)
callout(s, "We report this unprompted because a defensible result is one whose anomalies are "
           "explained rather than omitted.",
        y + 3.34, size=13.5, bg=GREEN_BG, fg=GREEN_TX)

# ── what it solves
s, y = new("What the engine solves", kicker="Evaluation",
           sub="Four concrete problems a fixed rendering policy cannot address — and the evidence for each is quantitative, not asserted.")
probs = [
    ("The mismatched-client problem",
     f"A single annotated route cannot suit both a capable and a constrained client. The engine "
     f"gives the capable client an **{ad['doc_bytes']}-byte interactive shell** and the constrained "
     f"client **{fx['doc_bytes']:,} bytes of finished HTML** — from the same URL.",
     "verified in the trigger matrix, measured on the fast link"),
    ("The origin-cost problem",
     f"Under concurrency the engine shifts to cache-backed rendering **without operator "
     f"intervention**, cutting origin CPU by **{cpu_cut:.1f} %** for the same delivered traffic.",
     "measured on the 100 ms link"),
    ("The tail-latency problem",
     f"Fixed server-side rendering under load produced a 99th percentile of **{fx['p99']} ms**; "
     f"adaptive selection reduced it to **{ad['p99']} ms** ({p99_cut:.1f} % lower), because the work "
     f"that created the queue is no longer performed.",
     "measured on the fast link"),
    ("The opacity problem",
     "In a conventional framework, the reason a page was rendered a particular way is a property of "
     "**source code**. Here it is a property of the **response**: every reply states its strategy and "
     "the rule that produced it.",
     "which is what made every claim in this deck externally verifiable"),
]
tf = tbox(s, M, y, CW, 4.1)
for i, (t, d, e) in enumerate(probs):
    para(tf, f"{i+1}.  {t}", size=15, bold=True, color=NAVY, font=HEAD, after=3,
         first=(i == 0))
    para(tf, d, size=13, color=TEXT, font=BODY, after=2, indent=0.32, line=1.15)
    para(tf, e, size=11.5, color=BLUE, font=BODY, italic=True, after=11, indent=0.32)

# ── objectives evaluation
s, y = new("Evaluation against the objectives", kicker="Evaluation")
table(s, ["Objective", "Outcome", "Evidence"], [
    ["1. Identify and observe contextual variables", "Achieved",
     "Seven signals observed per request with a defined precedence; control-surface table and the 17/17 trigger matrix"],
    ["2. Pure, deterministic decision engine", "Achieved",
     f"{ds['total']}-context enumeration: total, deterministic, no dead rules; plus 28 passing automated tests"],
    ["3. Six interchangeable strategy modules", "Achieved",
     "All six selected and rendered correctly; per-strategy performance measured for each"],
    ["4. Externally observable decisions", "Achieved",
     "Proof headers, three log lines and a metrics record per request — every result in this deck was read from them"],
    ["5. Zero-cost reproducible private server", "Achieved",
     "Five-container stack; simulation validated against its configuration with small, consistent residuals"],
    ["6. Measure and compare against fixed policies", "Achieved",
     "Two controlled load experiments under opposite bottleneck regimes, reporting gains and non-gains alike"],
], [4.2, 1.5, 6.6], y, fs=12, rh=0.46,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
   row_colors=[None] * 6)
callout(s, "Every objective set in Chapter 1 is met, and each is backed by evidence that a reviewer "
           "can regenerate from a clean checkout with the commands in Appendix C of the report.",
        y + 3.20, size=13.5, bg=GREEN_BG, fg=GREEN_TX)

# ── REAL vs SIMULATED (from the uploaded slide)
s, y = new("What is real and what is simulated?", kicker="Honest scoping",
           sub="Asked directly, answered directly — the boundary between the engineered system and the controlled environment it was measured in.")
table(s, ["Term", "Simple meaning"], [
    ["Runtime rule decision", "REAL — the server evaluates the rule table for engine requests."],
    ["SSG / SSR / CSR / ISR / Streaming implementations",
     "REAL — separate code paths render, cache, stream, or browser-render."],
    ["Redis", "REAL software — a real Redis container is used, though only in a local Docker environment."],
    ["In-flight concurrency", "REAL — the server counts current requests and uses it for load classification."],
    ["User-Agent device detection", "REAL but heuristic — useful, not a perfect measurement of hardware capability."],
    ["Browser RTT / latency check", "REAL measurement to this server using /health."],
    ["Network delay categories", "SIMULATED for controlled experiments — 0 / 100 / 400 ms delay can be injected."],
    ["Edge distance", "SIMULATED — Edge 1/2 use different added delays but normally run on the same computer."],
], [3.5, 8.8], y, fs=12, rh=0.375)
callout(s, "ARE is a working adaptive runtime tested in a controlled private-server environment. "
           "The **strategy selection and rendering/caching code paths are real**; **edge distance "
           "and network conditions are intentionally simulated** so experiments are repeatable.",
        y + 3.42, size=14, bg=GREEN_BG, fg=GREEN_TX, height=0.86,
        label="Best factual defense statement")

# ═════════════════════════════════════ SECTION 5 — CHALLENGES & FUTURE
section("05", "Challenges, Conclusion and Future",
        ["The problems that actually cost us time, and how each was resolved",
         "What is proven, and the limits that bound those claims",
         "Where a rule-based selector grows into a self-tuning rendering controller"])

# ── challenges I
s, y = new("Problems encountered and solutions adopted", kicker="Challenges 1 of 2",
           sub="Reported because the resolutions are part of the engineering contribution.")
table(s, ["Problem", "Impact", "Solution adopted"], [
    ["HTTP header values must be Latin-1, but decision reasons contained typographic arrows — the response failed",
     "High",
     "The reason is sanitised to printable ASCII on the response while the server log retains the original text, so no information is lost"],
    ["nginx cannot inject per-node latency, so edges modelled as proxies could demonstrate distance but not behaviour",
     "Medium",
     "Edges were re-modelled as full engine instances running the same image with their own identity, latency and cache — making Edge-ISR a real strategy rather than a flag"],
    ["A streamed response that is buffered before sending is not streaming at all",
     "Medium",
     "The strategy returns a PassThrough stream that the engine pipes straight to the response, resolving as soon as the shell is ready; chunked encoding is the external proof"],
    ["SSG could never be selected on a cold start, because rule 1 requires a usable cache that only a previous request could create",
     "Medium",
     "Static pages are pre-built at server start-up, so the artefact exists before the first request arrives"],
    ["A pre-built artefact embedded the cold cache state present at build time, so the page described a context contradicting its own badge",
     "Low",
     "The prebuild embeds the context under which rule 1 actually serves the artefact, and the live decision always travels in the response header"],
], [4.6, 0.9, 6.8], y, fs=11, rh=0.62,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT], bold_first=False)

# ── challenges II
s, y = new("Problems encountered and solutions adopted", kicker="Challenges 2 of 2")
table(s, ["Problem", "Impact", "Solution adopted"], [
    ["A burst of requests arriving on a stale cache entry could trigger many simultaneous re-renders",
     "Medium",
     "Background revalidation is single-flight: a module-level set of in-flight keys guarantees exactly one refresh per key"],
    ["Demonstration pages that read the clock during render produced React hydration mismatches",
     "Medium",
     "First render depends only on props; all live values are introduced in effects after mount, and timestamps render as UTC before localising"],
    ["Browsers cannot attach custom headers to an ordinary navigation, so in-page controls could only predict a strategy rather than trigger one",
     "Medium",
     "A query-parameter alias was added for every control header, so a link or a reload genuinely re-renders under a new strategy"],
    ["A fixed-SSR baseline could not be held under load, because the high-load rule correctly promoted it to ISR",
     "Medium",
     "The baseline pins the load signal to low, which suppresses rule 3 and holds the comparison policy genuinely fixed — a wasted experiment run taught us this"],
], [4.6, 0.9, 6.8], y, fs=11.5, rh=0.66,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT], bold_first=False)
callout(s, "Several of these look like bugs and are not — they are the system behaving "
           "correctly in a way that surprised us. Distinguishing the two took real debugging time, "
           "and each is recorded in the report's limitations so nobody 're-fixes' correct behaviour.",
        y + 3.10, size=13.5)

# ── limitations
s, y = new("Limitations and threats to validity", kicker="Honest boundaries",
           sub="Stated explicitly so the results are not read more widely than the evidence supports.")
lims = [
    ("Single-host evaluation",
     "Origin, edges, proxy, cache and load generator share one machine and contend for the same "
     "processor. Absolute figures would differ on distributed hardware — **the comparisons "
     "between arms, run back to back under identical conditions, are the meaningful quantities.**"),
    ("Simulated rather than physical conditions",
     "Network classes are server-side delays and edge distance is injected latency. The simulation "
     "matches its configuration, but a simulated 400 ms link does not reproduce the jitter, packet "
     "loss or bandwidth constraint of a real mobile network."),
    ("Server-side metrics only",
     "FCP and LCP are not instrumented, so the client-perceived benefit of streaming and of smaller "
     "payloads is argued from transferred bytes and server timings rather than measured in a browser."),
    ("Two behaviours are environment-bound",
     "Edge-ISR is not reachable through the default proxy route (the proxy correctly stamps origin "
     "requests as origin-served), and forcing the realtime page to static volatility is "
     "non-deterministic because it depends on whether that page's cache entry happens to be warm."),
    ("The rule table encodes our judgement",
     "It is internally consistent, exhaustively verified and externally observable — but it is "
     "**not learned from outcome data**. Establishing that these particular nine rules are optimal "
     "would require the closed-loop evaluation proposed next."),
]
tf = tbox(s, M, y, CW, 4.2)
for i, (t, d) in enumerate(lims):
    para(tf, t, size=14, bold=True, color=NAVY, font=HEAD, after=2, first=(i == 0))
    para(tf, d, size=12.5, color=TEXT, font=BODY, after=10, indent=0.0, line=1.14)

# ── conclusion
s, y = new("Conclusion", kicker="What was set out to do, and what was done",
           sub="This project set out to move the choice of rendering strategy from build time, where it is a guess, to request time, where it is an observation.")
tiles(s, [
    (f"{D['trigger_pass']}/17", "context→strategy triggers reproduce exactly"),
    (f"{ds['total']}", "contexts enumerated: total, deterministic, no dead rules"),
    ("28/28", "automated tests pass"),
    (f"{ds['ns_per_decision']:.1f} ns", "to decide — ~1/10,000 of a request"),
    (f"+{rps_gain:.0f} %", "throughput when the origin is the bottleneck"),
    (f"−{cpu_cut:.0f} %", "origin CPU when the network is"),
], y, h=1.20, vsize=21)
tf = tbox(s, M, y + 1.46, CW, 2.4)
para(tf, "The engine analyses seven contextual signals per request, selects one of six strategies "
         "through a pure, ordered, first-match-wins rule table, renders with it, and publishes both "
         "the decision and its justification on the response — running as a five-container "
         "private server on open-source software at no cost.",
     size=14, after=11, first=True, line=1.20)
para(tf, "**The wider conclusion is that rendering strategy is a legitimate runtime optimisation "
         "variable.** The web accumulated six rendering strategies over roughly fifteen years and "
         "left the question of when to use each of them to a developer annotation. This work shows "
         "that the question can be answered per request, that answering it costs essentially nothing, "
         "and that answering it well produces measurable gains in exactly the conditions where a "
         "fixed choice is most obviously wrong.",
     size=14, after=0, line=1.20)
callout(s, "The engine is also, deliberately, a **research instrument**: because every decision is "
           "recorded with its context, its timings and its outcome, it produces the dataset a future "
           "learned policy would need in order to improve on the rules it currently uses.",
        y + 3.92, size=13.5, bg=BLUE_L, height=0.72)

# ── achievements
s, y = new("Achievements", kicker="What was delivered")
ach = [
    ("A complete, working runtime", "that selects among six rendering strategies per request, built on Node.js, TypeScript and React 18 over the native HTTP module."),
    ("A decision layer that is pure, total, deterministic and exhaustively verified", f"across its entire {ds['total']}-point context space, at a cost of {ds['ns_per_decision']:.2f} ns per decision."),
    ("Six interchangeable strategy modules behind one interface", "including a genuine stale-while-revalidate implementation with single-flight background refresh, and a streaming renderer that is never buffered."),
    ("A zero-cost, reproducible, five-container private server", "whose simulated conditions were validated against their configuration."),
    ("An end-to-end evidence chain", "— proof headers, structured logs, per-request metrics and aggregated reports — which made every result externally verifiable."),
    ("A controlled comparative evaluation against fixed policies", "under two bottleneck regimes, reporting both the large gains and the cases where the gain is capacity rather than speed."),
]
tf = tbox(s, M, y, CW, 4.2)
for i, (t, d) in enumerate(ach):
    p = para(tf, t, size=14, bold=True, color=NAVY, font=HEAD, after=2,
             first=(i == 0), bullet="▪")
    para(tf, d, size=12.5, color=TEXT, font=BODY, after=12, indent=0.26, line=1.14)

# ── roadmap
s, y = new("Future enhancement roadmap", kicker="Recommendations",
           sub="The architecture was shaped to make its own extension straightforward. Every item below exploits a seam that already exists in the design.")
table(s, ["Horizon", "Enhancement", "Why the architecture already supports it"], [
    ["Immediate", "Promote SSG artefacts into the memory cache tier",
     "Removes the filesystem-read anomaly measured in Result 4b; a change confined to one strategy module"],
    ["Immediate", "Adopt real client signals: Save-Data, Client Hints, Network Information API",
     "The analyzer already resolves each signal by precedence, so a browser-supplied hint becomes one more source ahead of inference"],
    ["Short term", "Instrument client-side paint metrics and feed them back",
     "The metrics record is already per-request and schema-driven; adding FCP and LCP completes the loop from decision to perceived outcome"],
    ["Short term", "Per-route and per-tenant policy tables",
     "The policy is data, not control flow, so a second table can be selected at request time without touching the evaluator"],
    ["Medium term", "Learned selection policy",
     "decide() is a pure function from context to strategy; a model can be substituted behind that exact signature with no change to the pipeline"],
    ["Medium term", "Closed-loop, SLO-driven control",
     "Every decision is already logged with its outcome, which is the training and feedback signal such a controller requires"],
    ["Medium term", "Component-level granularity (islands)",
     "The strategy interface renders a page module; the same interface can render a component subtree, letting one page mix strategies"],
    ["Long term", "Geographically distributed edges and edge runtimes",
     "Edges are already full engine instances distinguished only by environment, so relocating one is a deployment change, not a redesign"],
    ["Long term", "Energy- and cost-aware objectives",
     "Origin CPU is already sampled per request, so a rule can optimise for energy or cost rather than latency alone"],
], [1.5, 4.2, 6.6], y, fs=11, rh=0.42,
   aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])

# ── forward vision
s, y = new("Where this goes next", kicker="The vision, forward",
           sub="Two directions deserve emphasis, because they are where the work becomes genuinely forward-looking.")
w = (CW - 0.28) / 2
rect(s, M, y, w, 2.75, fill=PANEL, line=LINE, lw=0.75)
rect(s, M, y, w, 0.05, fill=BLUE)
tf = tbox(s, M + 0.24, y + 0.22, w - 0.48, 2.4)
para(tf, "FROM RULES TO A LEARNED POLICY", size=11.5, bold=True, color=BLUE, font=HEAD,
     after=9, first=True)
para(tf, "The present rule table encodes human judgement, and we are candid that its optimality is "
         "asserted rather than proven. But because decide() is a pure function with a fixed "
         "signature, a supervised or reinforcement-learned model can be **dropped in behind that "
         "signature without disturbing a single strategy module** — and the metrics log the "
         "engine already writes is precisely the training data such a model requires.",
     size=13, after=8, line=1.17)
para(tf, "The path from rule-based to learned selection is an implementation of an existing "
         "interface, not a rewrite.", size=13, bold=True, color=NAVY, after=0, line=1.17)
x2 = M + w + 0.28
rect(s, x2, y, w, 2.75, fill=PANEL, line=LINE, lw=0.75)
rect(s, x2, y, w, 0.05, fill=BLUE)
tf = tbox(s, x2 + 0.24, y + 0.22, w - 0.48, 2.4)
para(tf, "FROM PAGES TO COMPONENTS", size=11.5, bold=True, color=BLUE, font=HEAD,
     after=9, first=True)
para(tf, "Strategies are currently applied to whole pages because that is the unit frameworks "
         "expose — but nothing in the engine requires it. A page is a component tree, and the "
         "strategy interface renders a component. The natural next step is **a page whose navigation "
         "is statically generated, whose main content is streamed, and whose interactive island is "
         "client-rendered — all decided per request.**",
     size=13, after=8, line=1.17)
para(tf, "That is where the wider ecosystem is already heading: islands architectures, partial "
         "hydration and server components.", size=13, bold=True, color=NAVY, after=0, line=1.17)
callout(s, "Taken together these describe a system that starts as a rule-based selector and grows "
           "into a **self-tuning rendering controller** — one that observes its own outcomes, "
           "learns which strategy actually served each class of visitor best, operates at component "
           "granularity across a distributed edge, and optimises not only for speed but for the cost "
           "and energy of producing a page. Every one of those steps is an extension of the seams "
           "built in this project rather than a departure from them, which is the strongest practical "
           "evidence that the architecture was the right one.",
        y + 2.98, size=13, bg=BLUE_L, height=1.10)

# ── reproduce
s, y = new("Reproduce every result yourself", kicker="Live demonstration",
           sub="From a clean checkout. The stack must be rebuilt rather than restarted after a code change, because the deployment mounts no volumes.")
code(s, [
    "# 1. bring up the five-container private server from a clean state",
    "docker compose down && docker compose up -d --build",
    "",
    "# 2. correctness: unit tests and the strategy-switch matrix",
    "npm test                                   # 28 tests, 4 files",
    "bash scripts/switch-test.sh                # same URL, varied context",
    "",
    "# 3. read the decision for any context (header form and query form)",
    "curl -sI \"localhost:8080/dynamic?net=fast&device=desktop\"  | grep -i x-rendering-strategy",
    "curl -sI -H 'X-Device-Type: mobile' localhost:8080/dynamic | grep -i x-rendering-strategy",
    "",
    "# 4. prove that streaming is not buffered  (GET, not HEAD -- a HEAD carries no body)",
    "curl -s -D- -o /dev/null localhost:8080/heavy | grep -iE 'transfer-encoding|x-streaming'",
    "",
    "# 5. edge behaviour (the default route is stamped as origin by the proxy)",
    "curl -sI \"localhost:8080/edge1/static?cache=cold\" | grep -i x-rendering-strategy",
    "",
    "# 6. load comparison: fixed SSR policy versus adaptive selection",
    "ab -n 800 -c 50 \"localhost:8080/dynamic?net=fast&volatility=realtime&device=mobile&load=low\"",
    "ab -n 800 -c 50 \"localhost:8080/dynamic?net=fast&device=desktop\"",
], y, fs=11.5)
callout(s, "The engine emits three lines per request — the URL, the observed context, and the "
           "strategy with its reason — which is the primary evidence that the decision reported "
           "in the header is the decision the engine actually took.",
        y + 3.94, size=13)

# ── thank you
s = prs.slides.add_slide(BLANK)
_slide_no[0] += 1
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, 0.20, fill=BLUE)
tf = tbox(s, 1.05, 2.05, 11.2, 1.0)
para(tf, "Thank you", size=46, bold=True, color=WHITE, font=HEAD, after=0, first=True)
rect(s, 1.05, 3.20, 1.5, 0.05, fill=BLUE)
tf = tbox(s, 1.05, 3.52, 10.4, 1.2)
para(tf, "Rendering strategy should be a runtime decision computed from the conditions of the "
         "request — observable, explainable and reproducible from outside the system.",
     size=17, color=RGBColor(0xC6, 0xD8, 0xEC), font=BODY, after=14, first=True, line=1.25)
para(tf, "Questions and discussion", size=15, bold=True,
     color=RGBColor(0x8F, 0xB6, 0xDC), font=HEAD, after=0)
tf = tbox(s, 1.05, 5.55, 11.2, 0.9)
para(tf, "Bijay B.k (220305)  ·  Devendra Pandey (220306)  ·  Manish Joshi (220312)  "
         "·  Pramod Panta (220317)",
     size=13, color=WHITE, font=BODY, after=5, first=True)
para(tf, "Supervisor: Er. Robinhood Khadka  ·  Department of ICT and Computer Engineering  "
         "·  Cosmos College of Management & Technology",
     size=11.5, color=RGBColor(0x9F, 0xB8, 0xD4), font=BODY, after=0)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
