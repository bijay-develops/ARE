#!/usr/bin/env python3
"""
Render the 20-slide Adaptive Rendering Engine defense deck.

One content model (deck_content.py) -> three outputs:
    ARE_Defense_Slides.pdf     slide-shaped pages, the reference rendering
    ARE_Defense_Slides.docx    same content, editable, one slide per page
    ARE_Defense_Slides.pptx    the presentable deck

PDF and PPTX share an absolute layout engine that measures every block and
auto-shrinks a slide until it fits, so no slide can ever overflow. DOCX is
emitted as flowed Word content (real headings, tables and bullets).

    python3 report-build/build_defense_deck.py
"""

import os
import re
import sys

from PIL import Image
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import deck_content as C  # noqa: E402

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_PDF = os.path.join(ROOT, "ARE_Defense_Slides.pdf")
OUT_DOCX = os.path.join(ROOT, "ARE_Defense_Slides.docx")
OUT_PPTX = os.path.join(ROOT, "ARE_Defense_Slides.pptx")

SUP = "/System/Library/Fonts/Supplemental"

# ────────────────────────────────────────────────────────────── design tokens
PW, PH = 13.333, 7.5
MG = 0.62
CW = PW - 2 * MG

NAVY = (0x1F, 0x38, 0x64)
NAVY_D = (0x14, 0x25, 0x44)
BLUE = (0x1F, 0x6F, 0xB2)
BLUE_PALE = (0xE7, 0xEF, 0xF8)
BLUE_TINT = (0xC9, 0xDB, 0xEE)
ROW_A = (0xFF, 0xFF, 0xFF)
ROW_B = (0xEF, 0xF3, 0xFA)
PANEL = (0xF5, 0xF7, 0xFB)
CODE_BG = (0xF2, 0xF5, 0xF9)
GREEN_BG = (0xE5, 0xF1, 0xDD)
GREEN_TX = (0x2B, 0x58, 0x1D)
AMBER_BG = (0xFD, 0xF3, 0xDD)
AMBER_TX = (0x78, 0x50, 0x00)
TEXT = (0x23, 0x28, 0x32)
MUTED = (0x5B, 0x65, 0x75)
WHITE = (0xFF, 0xFF, 0xFF)
LINE = (0xC6, 0xD3, 0xE4)

# Run style is a bitmask: 0 normal, 1 bold, 2 italic, 3 bold-italic.
# (True == 1 in Python, so plain `(text, True)` literals still mean bold.)
BOLD, ITAL = 1, 2

# family -> (office name, {style: (pdf font name, file)})
FAMS = {
    "body": ("Georgia", {
        0: ("Body", f"{SUP}/Georgia.ttf"),
        1: ("Body-B", f"{SUP}/Georgia Bold.ttf"),
        2: ("Body-I", f"{SUP}/Georgia Italic.ttf"),
        3: ("Body-BI", f"{SUP}/Georgia Bold Italic.ttf")}),
    "head": ("Trebuchet MS", {
        0: ("Head", f"{SUP}/Trebuchet MS.ttf"),
        1: ("Head-B", f"{SUP}/Trebuchet MS Bold.ttf"),
        2: ("Head-I", f"{SUP}/Trebuchet MS Italic.ttf"),
        3: ("Head-BI", f"{SUP}/Trebuchet MS Bold Italic.ttf")}),
    "mono": ("Courier New", {
        0: ("Mono", f"{SUP}/Courier New.ttf"),
        1: ("Mono-B", f"{SUP}/Courier New Bold.ttf"),
        2: ("Mono-I", f"{SUP}/Courier New Italic.ttf"),
        3: ("Mono-BI", f"{SUP}/Courier New Bold Italic.ttf")}),
}
for _fam, (_office, _styles) in FAMS.items():
    for _st, (_name, _path) in _styles.items():
        pdfmetrics.registerFont(TTFont(_name, _path))

BASE = 11.0          # body size at scale 1.0
LEAD = 1.20          # line height multiplier


def pdf_font(fam, style):
    return FAMS[fam][1][int(style) & 3][0]


def sw(text, fam, style, size):
    return pdfmetrics.stringWidth(text, pdf_font(fam, style), size) / 72.0


# ─────────────────────────────────────────────────────────── rich-text runs
_EMPH = re.compile(r"\*\*(.+?)\*\*|\*(.+?)\*", re.S)


def parse_runs(text):
    """'a **b** and *c*' -> [('a ',0), ('b',BOLD), (' and ',0), ('c',ITAL)]"""
    text = str(text)
    out, pos = [], 0
    for m in _EMPH.finditer(text):
        if m.start() > pos:
            out.append((text[pos:m.start()], 0))
        if m.group(1) is not None:
            out.append((m.group(1), BOLD))
        else:
            out.append((m.group(2), ITAL))
        pos = m.end()
    if pos < len(text):
        out.append((text[pos:], 0))
    return out or [("", 0)]


def wrap_runs(runs, fam, size, maxw):
    """Greedy word wrap across styled runs -> list of lines, each [(text,bold)].

    Words are tokenised from the *concatenated* text, so a word whose weight
    changes mid-way (e.g. '**arrived**.') stays one word and never acquires a
    spurious space before its punctuation.
    """
    words, cur = [], []          # each word: [(text, bold), ...]
    for txt, bold in runs:
        parts = txt.split(" ")
        for i, part in enumerate(parts):
            if i > 0:
                if cur:
                    words.append(cur)
                cur = []
            if part:
                if cur and cur[-1][1] == bold:
                    cur[-1] = (cur[-1][0] + part, bold)
                else:
                    cur.append((part, bold))
    if cur:
        words.append(cur)
    if not words:
        return [[("", False)]]

    def wwidth(word):
        return sum(sw(t, fam, bd, size) for t, bd in word)

    def hard_split(word):
        """Break an unbreakable token (a long URL, STREAMING_SSR) at character
        level so it can never spill outside its column."""
        chunks, chunk, cw = [], [], 0.0
        for t, bd in word:
            for ch in t:
                w1 = sw(ch, fam, bd, size)
                if chunk and cw + w1 > maxw:
                    chunks.append(chunk)
                    chunk, cw = [], 0.0
                if chunk and chunk[-1][1] == bd:
                    chunk[-1] = (chunk[-1][0] + ch, bd)
                else:
                    chunk.append((ch, bd))
                cw += w1
        if chunk:
            chunks.append(chunk)
        return chunks

    # (segments, glued) — glued means "no space before me", used for split tokens
    tokens = []
    for word in words:
        if wwidth(word) > maxw:
            for i, part in enumerate(hard_split(word)):
                tokens.append((part, i > 0))
        else:
            tokens.append((word, False))

    space = sw(" ", fam, False, size)
    lines, line, linew = [], [], 0.0
    for word, glued in tokens:
        ww = wwidth(word)
        add = ww if (not line or glued) else ww + space
        if line and linew + add > maxw:
            lines.append(line)
            line, linew = [(word, glued)], ww
        else:
            line.append((word, glued))
            linew += add
    if line:
        lines.append(line)

    out = []
    for ln in lines:
        segs = []
        for wi, (word, glued) in enumerate(ln):
            if wi > 0 and segs and not glued:        # the joining space
                segs[-1] = (segs[-1][0] + " ", segs[-1][1])
            for t, bd in word:
                if segs and segs[-1][1] == bd:
                    segs[-1] = (segs[-1][0] + t, bd)
                else:
                    segs.append((t, bd))
        out.append(segs or [("", False)])
    return out


def line_h(size):
    return size / 72.0 * LEAD


# ═══════════════════════════════════════════════════════════ layout engine
class L:
    """Accumulates absolute primitives for one slide."""

    def __init__(self, s):
        self.p = []
        self.s = s          # scale factor

    def rect(self, x, y, w, h, fill=None, stroke=None, lw=0.6):
        self.p.append(("rect", x, y, w, h, fill, stroke, lw))

    def line(self, x, y, w, size, runs, color, fam="body", align="l", bold=False):
        self.p.append(("text", x, y, w, size, runs, color, fam, align, bold))

    def img(self, x, y, w, h, path):
        self.p.append(("img", x, y, w, h, path))


def txt_block(lo, x, y, w, text, size, color, fam="body", gap=0.0, bold=False,
              align="l"):
    """Wrapped paragraph. Returns new y."""
    runs = parse_runs(text)
    if bold:
        runs = [(t, True) for t, _ in runs]
    for ln in wrap_runs(runs, fam, size, w):
        lo.line(x, y, w, size, ln, color, fam, align, bold)
        y += line_h(size)
    return y + gap


def measure_txt(text, w, size, fam="body", gap=0.0):
    return len(wrap_runs(parse_runs(text), fam, size, w)) * line_h(size) + gap


# ── individual blocks ───────────────────────────────────────────────────────
def blk(lo, b, x, y, w, measure=False):
    """Lay out one block at (x, y) with width w. Returns the bottom y."""
    s = lo.s
    kind = b[0]

    if kind == "gap":
        return y + b[1] * s

    if kind == "head":
        size = 9.5 * s
        y = txt_block(lo, x, y, w, b[1].upper(), size, BLUE, "head", bold=True)
        return y + 0.07 * s

    if kind == "bul":
        size = BASE * s
        sub_size = (BASE - 1.2) * s
        ind = 0.22 * s
        for it in b[1]:
            sub = None
            if isinstance(it, tuple):
                it, sub = it
            runs = parse_runs(it)
            lines = wrap_runs(runs, "body", size, w - ind)
            for i, ln in enumerate(lines):
                if i == 0:
                    lo.line(x, y, ind, size, [("▪", False)], BLUE, "body")
                lo.line(x + ind, y, w - ind, size, ln, TEXT, "body")
                y += line_h(size)
            y += 0.035 * s
            if sub:
                for ln in wrap_runs(parse_runs(sub), "body", sub_size, w - ind):
                    lo.line(x + ind, y, w - ind, sub_size, ln, MUTED, "body")
                    y += line_h(sub_size)
                y += 0.045 * s
            y += 0.055 * s
        return y

    if kind == "kv":
        tsize = (BASE - 0.3) * s
        dsize = (BASE - 1.0) * s
        for term, desc in b[1]:
            for ln in wrap_runs(parse_runs(term), "head", tsize, w):
                lo.line(x, y, w, tsize, [(t, True) for t, _ in ln], NAVY, "head")
                y += line_h(tsize)
            y += 0.015 * s
            for ln in wrap_runs(parse_runs(desc), "body", dsize, w - 0.10):
                lo.line(x + 0.10, y, w - 0.10, dsize, ln, TEXT, "body")
                y += line_h(dsize)
            y += 0.105 * s
        return y

    if kind == "code":
        lines, cap = b[1], (b[2] if len(b) > 2 else None)
        size = 9.0 * s
        lh = size / 72.0 * 1.34
        pad = 0.11 * s
        caph = (0.20 * s if cap else 0.0)
        h = pad * 2 + caph + len(lines) * lh
        lo.rect(x, y, w, h, CODE_BG, LINE, 0.6)
        yy = y + pad
        if cap:
            lo.line(x + 0.13, yy, w - 0.26, 8.5 * s, [(cap, True)], BLUE, "head")
            yy += caph
        for ln in lines:
            col = MUTED if ln.strip().startswith(("#", "//")) else NAVY_D
            lo.line(x + 0.13, yy, w - 0.26, size, [(ln, False)], col, "mono")
            yy += lh
        return y + h + 0.06 * s

    if kind == "note":
        text, kindn = b[1], (b[2] if len(b) > 2 else "info")
        bg, fg, bar = {
            "info": (BLUE_PALE, NAVY, BLUE),
            "good": (GREEN_BG, GREEN_TX, GREEN_TX),
            "warn": (AMBER_BG, AMBER_TX, AMBER_TX),
            "dark": (NAVY, WHITE, BLUE),
        }[kindn]
        size = (BASE + 0.3) * s
        pad = 0.12 * s
        inner = w - 0.42 * s
        lines = wrap_runs(parse_runs(text), "body", size, inner)
        h = pad * 2 + len(lines) * line_h(size)
        lo.rect(x, y, w, h, bg)
        lo.rect(x, y, 0.055, h, bar)
        yy = y + pad
        for ln in lines:
            lo.line(x + 0.22 * s, yy, inner, size, ln, fg, "body")
            yy += line_h(size)
        return y + h + 0.09 * s

    if kind == "img":
        path, maxh, cap = b[1], b[2], (b[3] if len(b) > 3 else None)
        if not path or not os.path.exists(path):
            return y
        iw, ih = Image.open(path).size
        ar = iw / ih
        h = maxh * s
        ww = h * ar
        if ww > w:
            ww = w
            h = ww / ar
        cx = x + (w - ww) / 2
        lo.rect(cx - 0.025, y - 0.025, ww + 0.05, h + 0.05, WHITE, LINE, 0.6)
        lo.img(cx, y, ww, h, path)
        y += h + 0.07 * s
        if cap:
            size = 8.6 * s
            for ln in wrap_runs(parse_runs(cap), "body", size, w - 0.4):
                lo.line(x + 0.2, y, w - 0.4, size, ln, MUTED, "body", align="c")
                y += line_h(size)
            y += 0.05 * s
        return y

    if kind == "tiles":
        items = b[1]
        n = len(items)
        gap = 0.12 * s
        tw = (w - gap * (n - 1)) / n
        h = 0.86 * s
        for i, (val, lab) in enumerate(items):
            tx = x + i * (tw + gap)
            lo.rect(tx, y, tw, h, PANEL, LINE, 0.6)
            lo.rect(tx, y, tw, 0.042, BLUE)
            lo.line(tx, y + 0.14 * s, tw, 17 * s, [(val, True)], NAVY, "head", align="c")
            lab_s = 8.8 * s
            yy = y + 0.14 * s + line_h(17 * s) + 0.02
            for ln in wrap_runs(parse_runs(lab), "body", lab_s, tw - 0.14):
                lo.line(tx + 0.07, yy, tw - 0.14, lab_s, ln, TEXT, "body", align="c")
                yy += line_h(lab_s)
        return y + h + 0.10 * s

    if kind == "table":
        headers, rows, widths, opts = b[1], b[2], b[3], (b[4] if len(b) > 4 else {})
        fs = opts.get("fs", 10.2) * s
        hs = opts.get("hs", fs)
        center = set(opts.get("center", []))
        mono = set(opts.get("mono", []))
        rules = set(opts.get("rules", []))
        pad = 0.055 * s
        tot = sum(widths)
        wd = [q * (w / tot) for q in widths]

        def row_h(cells, size, fams):
            m = 0.0
            for i, cell in enumerate(cells):
                lns = wrap_runs(parse_runs(cell), fams[i], size, wd[i] - 2 * pad - 0.04)
                m = max(m, len(lns) * line_h(size))
            return m + 2 * pad

        hf = ["head"] * len(headers)
        hh = row_h(headers, hs, hf)
        lo.rect(x, y, w, hh, NAVY)
        cx = x
        for i, htxt in enumerate(headers):
            al = "c" if i in center else "l"
            lo.line(cx + pad, y + pad, wd[i] - 2 * pad, hs, [(htxt, True)], WHITE, "head", al)
            cx += wd[i]
        yy = y + hh
        for ri, row in enumerate(rows):
            fams = ["mono" if i in mono else "body" for i in range(len(row))]
            rh = row_h(row, fs, fams)
            lo.rect(x, yy, w, rh, ROW_A if ri % 2 == 0 else ROW_B)
            if ri in rules:
                lo.rect(x, yy, w, 0.012, BLUE_TINT)
            cx = x
            for i, cell in enumerate(row):
                al = "c" if i in center else "l"
                first = (i == 0)
                col = NAVY if first else TEXT
                lns = wrap_runs(parse_runs(cell), fams[i], fs, wd[i] - 2 * pad - 0.04)
                ty = yy + pad
                for ln in lns:
                    if first:
                        ln = [(t, True) for t, _ in ln]
                    lo.line(cx + pad, ty, wd[i] - 2 * pad, fs, ln, col, fams[i], al)
                    ty += line_h(fs)
                cx += wd[i]
            lo.rect(x, yy + rh - 0.006, w, 0.006, LINE)
            yy += rh
        return yy + 0.10 * s

    if kind == "split":
        left, right, frac = b[1], b[2], b[3]
        gap = 0.30 * s
        lw = (w - gap) * frac
        rw = w - gap - lw
        yl = y
        for sb in left:
            yl = blk(lo, sb, x, yl, lw)
        yr = y
        for sb in right:
            yr = blk(lo, sb, x + lw + gap, yr, rw)
        return max(yl, yr)

    raise ValueError(f"unknown block {kind}")


def layout_slide(sl):
    """Auto-shrink until the slide fits. Returns (primitives, scale)."""
    # grow to fill the slide, then shrink if necessary — first fitting scale wins
    ladder = [1.22, 1.18, 1.14, 1.10, 1.06, 1.03, 1.00, 0.97, 0.94, 0.91,
              0.88, 0.85, 0.82, 0.79, 0.76, 0.73, 0.70]
    for s in ladder:
        lo = L(s)
        y = 0.40
        if sl.get("kicker"):
            lo.line(MG, y, CW, 9.6, [(sl["kicker"].upper(), True)], BLUE, "head")
            y += 0.26
        tsize = 21.5
        for ln in wrap_runs(parse_runs(sl["title"]), "head", tsize, CW):
            lo.line(MG, y, CW, tsize, [(t, True) for t, _ in ln], NAVY, "head")
            y += line_h(tsize)
        y += 0.045
        if sl.get("sub"):
            y = txt_block(lo, MG, y, CW, sl["sub"], 11.2, MUTED, "body", gap=0.04)
        lo.rect(MG, y + 0.06, 1.05, 0.042, BLUE)
        y += 0.30
        top = y
        for b in sl["blocks"]:
            y = blk(lo, b, MG, y, CW)
        if y <= PH - 0.55 or s <= 0.70:
            return lo, s, (y, top)
    return lo, s, (y, top)


# ═════════════════════════════════════════════════════════════════════ PDF
def rgb(c):
    return (c[0] / 255.0, c[1] / 255.0, c[2] / 255.0)


def render_pdf(path):
    cv = rl_canvas.Canvas(path, pagesize=(PW * 72, PH * 72))
    cv.setTitle("Adaptive Rendering Engine — Final Defense Presentation")
    cv.setAuthor(", ".join(n for n, _ in C.TEAM))

    def Y(y):
        return (PH - y) * 72

    # ── title page
    cv.setFillColorRGB(*rgb(NAVY))
    cv.rect(0, 0, PW * 72, PH * 72, stroke=0, fill=1)
    cv.setFillColorRGB(*rgb(BLUE))
    cv.rect(0, Y(0.20), PW * 72, 0.20 * 72, stroke=0, fill=1)
    cv.setFillColorRGB(*rgb((0x8F, 0xB6, 0xDC)))
    cv.setFont("Head-B", 11.5)
    cv.drawString(MG * 72 + 30, Y(1.55), "FINAL YEAR ENGINEERING PROJECT  ·  DEFENSE PRESENTATION")
    cv.setFillColorRGB(*rgb(WHITE))
    cv.setFont("Head-B", 42)
    cv.drawString(MG * 72 + 30, Y(2.25), C.TITLE)
    cv.setFillColorRGB(*rgb((0xC6, 0xD8, 0xEC)))
    cv.setFont("Body", 14.5)
    yy = 2.62
    for ln in wrap_runs(parse_runs(C.SUBTITLE), "body", 14.5, 9.6):
        cv.drawString(MG * 72 + 30, Y(yy + 0.22), "".join(t for t, _ in ln))
        yy += line_h(14.5)
    cv.setFillColorRGB(*rgb(BLUE))
    cv.rect(MG * 72 + 30, Y(yy + 0.52), 1.35 * 72, 4, stroke=0, fill=1)
    cv.setFillColorRGB(*rgb((0x8F, 0xB6, 0xDC)))
    cv.setFont("Head-B", 10)
    cv.drawString(MG * 72 + 30, Y(4.30), "PRESENTED BY")
    cv.drawString(7.30 * 72, Y(4.30), "UNDER THE SUPERVISION OF")
    cv.setFillColorRGB(*rgb(WHITE))
    cv.setFont("Body", 13)
    ty = 4.60
    for nm, rl in C.TEAM:
        cv.drawString(MG * 72 + 30, Y(ty), f"{nm}   —   {rl}")
        ty += 0.28
    cv.setFont("Body", 13)
    cv.drawString(7.30 * 72, Y(4.60), C.SUPERVISOR)
    cv.setFillColorRGB(*rgb((0xC6, 0xD8, 0xEC)))
    cv.setFont("Body", 11.5)
    cv.drawString(7.30 * 72, Y(5.02), C.DEPT)
    cv.drawString(7.30 * 72, Y(5.28), C.COLLEGE)
    cv.setFillColorRGB(*rgb((0x9F, 0xB8, 0xD4)))
    cv.drawString(7.30 * 72, Y(5.54), C.DATE)
    cv.showPage()

    n = 1
    for sl in C.SLIDES:
        if sl.get("kind") == "title":
            continue
        n += 1
        lo, s, _ = layout_slide(sl)
        cv.setFillColorRGB(1, 1, 1)
        cv.rect(0, 0, PW * 72, PH * 72, stroke=0, fill=1)
        for p in lo.p:
            if p[0] == "rect":
                _, x, y, w, h, fill, stroke, lw = p
                if fill:
                    cv.setFillColorRGB(*rgb(fill))
                if stroke:
                    cv.setStrokeColorRGB(*rgb(stroke))
                    cv.setLineWidth(lw)
                cv.rect(x * 72, Y(y + h), w * 72, h * 72,
                        stroke=1 if stroke else 0, fill=1 if fill else 0)
            elif p[0] == "img":
                _, x, y, w, h, path_ = p
                cv.drawImage(path_, x * 72, Y(y + h), w * 72, h * 72,
                             preserveAspectRatio=True, anchor="c", mask="auto")
            else:
                _, x, y, w, size, runs, color, fam, align, _b = p
                cv.setFillColorRGB(*rgb(color))
                total = sum(sw(t, fam, bd, size) for t, bd in runs)
                if align == "c":
                    cx = x + (w - total) / 2
                elif align == "r":
                    cx = x + w - total
                else:
                    cx = x
                base = Y(y + size / 72.0 * 0.98)
                for t, bd in runs:
                    cv.setFont(pdf_font(fam, bd), size)
                    cv.drawString(cx * 72, base, t)
                    cx += sw(t, fam, bd, size)
        # footer
        cv.setFillColorRGB(*rgb(MUTED))
        cv.setFont("Head", 8.5)
        cv.drawString(MG * 72, Y(PH - 0.30), f"{C.TITLE}  ·  Final Defense")
        cv.drawRightString((PW - MG) * 72, Y(PH - 0.30), f"{n} / {len(C.SLIDES)}")
        cv.showPage()
    cv.save()
    return n


# ════════════════════════════════════════════════════════════════════ PPTX
def render_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(PW)
    prs.slide_height = Inches(PH)
    blank = prs.slide_layouts[6]

    def col(c):
        return RGBColor(*c)

    def add_rect(sl, x, y, w, h, fill, stroke=None, lw=0.6):
        sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = col(fill)
        else:
            sh.fill.background()
        if stroke:
            sh.line.color.rgb = col(stroke)
            sh.line.width = Pt(lw)
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def add_line(sl, x, y, w, size, runs, color, fam, align):
        h = size / 72.0 * 1.6
        box = sl.shapes.add_textbox(Inches(x), Inches(y - 0.035), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[align]
        for t, bd in runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.name = FAMS[fam][0]
            r.font.bold = bool(int(bd) & BOLD)
            r.font.italic = bool(int(bd) & ITAL)
            r.font.color.rgb = col(color)

    # title slide
    sl = prs.slides.add_slide(blank)
    add_rect(sl, 0, 0, PW, PH, NAVY)
    add_rect(sl, 0, 0, PW, 0.20, BLUE)
    add_line(sl, MG + 0.42, 1.42, 11, 11.5,
             [("FINAL YEAR ENGINEERING PROJECT  ·  DEFENSE PRESENTATION", True)],
             (0x8F, 0xB6, 0xDC), "head", "l")
    add_line(sl, MG + 0.42, 1.88, 11.5, 42, [(C.TITLE, True)], WHITE, "head", "l")
    yy = 2.72
    for ln in wrap_runs(parse_runs(C.SUBTITLE), "body", 14.5, 9.6):
        add_line(sl, MG + 0.42, yy, 9.8, 14.5, ln, (0xC6, 0xD8, 0xEC), "body", "l")
        yy += line_h(14.5)
    add_rect(sl, MG + 0.42, yy + 0.30, 1.35, 0.05, BLUE)
    add_line(sl, MG + 0.42, 4.24, 4, 10, [("PRESENTED BY", True)], (0x8F, 0xB6, 0xDC), "head", "l")
    add_line(sl, 7.30, 4.24, 5, 10, [("UNDER THE SUPERVISION OF", True)],
             (0x8F, 0xB6, 0xDC), "head", "l")
    ty = 4.56
    for nm, rl in C.TEAM:
        add_line(sl, MG + 0.42, ty, 5, 13, [(f"{nm}   —   {rl}", False)], WHITE, "body", "l")
        ty += 0.28
    add_line(sl, 7.30, 4.56, 5.4, 13, [(C.SUPERVISOR, False)], WHITE, "body", "l")
    add_line(sl, 7.30, 4.98, 5.4, 11.5, [(C.DEPT, False)], (0xC6, 0xD8, 0xEC), "body", "l")
    add_line(sl, 7.30, 5.24, 5.4, 11.5, [(C.COLLEGE, False)], (0xC6, 0xD8, 0xEC), "body", "l")
    add_line(sl, 7.30, 5.50, 5.4, 11.5, [(C.DATE, False)], (0x9F, 0xB8, 0xD4), "body", "l")

    n = 1
    for s_ in C.SLIDES:
        if s_.get("kind") == "title":
            continue
        n += 1
        lo, _, _ = layout_slide(s_)
        sl = prs.slides.add_slide(blank)
        for p in lo.p:
            if p[0] == "rect":
                add_rect(sl, p[1], p[2], p[3], p[4], p[5], p[6], p[7])
            elif p[0] == "img":
                sl.shapes.add_picture(p[5], Inches(p[1]), Inches(p[2]),
                                      Inches(p[3]), Inches(p[4]))
            else:
                add_line(sl, p[1], p[2], p[3], p[4], p[5], p[6], p[7], p[8])
        add_line(sl, MG, PH - 0.36, 6, 8.5, [(f"{C.TITLE}  ·  Final Defense", False)],
                 MUTED, "head", "l")
        add_line(sl, PW - MG - 1.2, PH - 0.36, 1.2, 8.5,
                 [(f"{n} / {len(C.SLIDES)}", False)], MUTED, "head", "r")
    prs.save(path)
    return n


# ════════════════════════════════════════════════════════════════════ DOCX
def render_docx(path):
    import docx
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.enum.section import WD_ORIENT
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = docx.Document()
    sec = doc.sections[0]
    sec.orientation = WD_ORIENT.LANDSCAPE
    sec.page_width = Inches(PW)
    sec.page_height = Inches(PH)
    sec.left_margin = sec.right_margin = Inches(MG)
    sec.top_margin = Inches(0.42)
    sec.bottom_margin = Inches(0.38)
    style = doc.styles["Normal"]
    style.font.name = "Georgia"
    style.font.size = Pt(10.5)
    style.paragraph_format.space_after = Pt(0)
    style.paragraph_format.line_spacing = 1.06

    def shade(cell, hexcolor):
        sh = OxmlElement("w:shd")
        sh.set(qn("w:val"), "clear")
        sh.set(qn("w:fill"), hexcolor)
        cell._tc.get_or_add_tcPr().append(sh)

    def nopad(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        mar = OxmlElement("w:tcMar")
        for side, v in (("top", 40), ("start", 60), ("bottom", 40), ("end", 60)):
            e = OxmlElement(f"w:{side}")
            e.set(qn("w:w"), str(v))
            e.set(qn("w:type"), "dxa")
            mar.append(e)
        tcPr.append(mar)

    def no_borders(table):
        tblPr = table._tbl.tblPr
        borders = OxmlElement("w:tblBorders")
        for edge in ("top", "left", "bottom", "right", "insideH", "insideV"):
            e = OxmlElement(f"w:{edge}")
            e.set(qn("w:val"), "none")
            e.set(qn("w:sz"), "0")
            borders.append(e)
        tblPr.append(borders)

    def hx(c):
        return "%02X%02X%02X" % c

    def para(ct, text, size=10.5, color=TEXT, fam="Georgia", bold=False, before=0,
             after=3, align=None, indent=0.0, bullet=False):
        p = ct.add_paragraph()
        pf = p.paragraph_format
        pf.space_before = Pt(before)
        pf.space_after = Pt(after)
        pf.line_spacing = 1.06
        if bullet:
            pf.left_indent = Inches(indent + 0.20)
            pf.first_line_indent = Inches(-0.20)
        elif indent:
            pf.left_indent = Inches(indent)
        if align:
            p.alignment = align
        if bullet:
            r = p.add_run("\u25aa  ")
            r.font.size = Pt(size)
            r.font.name = "Georgia"
            r.font.color.rgb = RGBColor(*BLUE)
        for t, bd in parse_runs(text):
            r = p.add_run(t)
            r.font.size = Pt(size)
            r.font.name = fam
            r.font.bold = bold or bool(int(bd) & BOLD)
            r.font.italic = bool(int(bd) & ITAL)
            r.font.color.rgb = RGBColor(*color)
        return p

    def clear_cell(cell):
        for p in list(cell.paragraphs):
            p._p.getparent().remove(p._p)

    def emit(ct, blocks, width=CW):
        for b in blocks:
            k = b[0]
            if k == "gap":
                para(ct, "", after=4)
            elif k == "head":
                para(ct, b[1].upper(), size=9, color=BLUE, fam="Trebuchet MS",
                     bold=True, before=6, after=3)
            elif k == "bul":
                for it in b[1]:
                    sub = None
                    if isinstance(it, tuple):
                        it, sub = it
                    para(ct, it, size=10.5, bullet=True, after=(1 if sub else 5))
                    if sub:
                        para(ct, sub, size=9.5, color=MUTED, indent=0.20, after=5)
            elif k == "kv":
                for term, desc in b[1]:
                    para(ct, term, size=10.5, color=NAVY, fam="Trebuchet MS",
                         bold=True, after=1, before=4)
                    para(ct, desc, size=10, indent=0.12, after=3)
            elif k == "code":
                if len(b) > 2 and b[2]:
                    para(ct, b[2], size=8.5, color=BLUE, fam="Trebuchet MS",
                         bold=True, after=2)
                t = ct.add_table(rows=1, cols=1)
                no_borders(t)
                cell = t.cell(0, 0)
                cell.width = Inches(width)
                shade(cell, hx(CODE_BG))
                nopad(cell)
                clear_cell(cell)
                for ln in b[1]:
                    p = cell.add_paragraph()
                    p.paragraph_format.space_after = Pt(0)
                    p.paragraph_format.line_spacing = 1.0
                    r = p.add_run(ln)
                    r.font.size = Pt(8.5)
                    r.font.name = "Courier New"
                    r.font.color.rgb = RGBColor(*NAVY_D)
                para(ct, "", after=3)
            elif k == "note":
                kindn = b[2] if len(b) > 2 else "info"
                bg, fg = {"info": (BLUE_PALE, NAVY), "good": (GREEN_BG, GREEN_TX),
                          "warn": (AMBER_BG, AMBER_TX), "dark": (NAVY, WHITE)}[kindn]
                t = ct.add_table(rows=1, cols=1)
                no_borders(t)
                cell = t.cell(0, 0)
                cell.width = Inches(width)
                shade(cell, hx(bg))
                clear_cell(cell)
                p = cell.add_paragraph()
                p.paragraph_format.space_before = Pt(3)
                p.paragraph_format.space_after = Pt(3)
                p.paragraph_format.line_spacing = 1.08
                for tx, bd in parse_runs(b[1]):
                    r = p.add_run(tx)
                    r.font.size = Pt(10.5)
                    r.font.name = "Georgia"
                    r.font.bold = bool(int(bd) & BOLD)
                    r.font.italic = bool(int(bd) & ITAL)
                    r.font.color.rgb = RGBColor(*fg)
                para(ct, "", after=3)
            elif k == "img":
                path_, maxh = b[1], b[2]
                if not path_ or not os.path.exists(path_):
                    continue
                iw, ih = Image.open(path_).size
                w = min(width, maxh * (iw / ih))
                p = ct.add_paragraph()
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                p.paragraph_format.space_after = Pt(4)
                p.add_run().add_picture(path_, width=Inches(w))
                if len(b) > 3 and b[3]:
                    para(ct, b[3], size=8.5, color=MUTED,
                         align=WD_ALIGN_PARAGRAPH.CENTER, after=5)
            elif k == "tiles":
                t = ct.add_table(rows=2, cols=len(b[1]))
                no_borders(t)
                for i, (val, lab) in enumerate(b[1]):
                    for ri, (txt, sz, cl, bd) in enumerate(
                            [(val, 15, NAVY, True), (lab, 8.5, TEXT, False)]):
                        c = t.cell(ri, i)
                        c.width = Inches(width / len(b[1]))
                        shade(c, hx(PANEL))
                        clear_cell(c)
                        p = c.add_paragraph()
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        p.paragraph_format.space_after = Pt(0)
                        r = p.add_run(txt)
                        r.font.size = Pt(sz)
                        r.font.name = "Trebuchet MS"
                        r.font.bold = bd
                        r.font.color.rgb = RGBColor(*cl)
                para(ct, "", after=4)
            elif k == "table":
                headers, rows, widths, opts = b[1], b[2], b[3], (b[4] if len(b) > 4 else {})
                fs = opts.get("fs", 10.2) - 0.6
                center = set(opts.get("center", []))
                mono = set(opts.get("mono", []))
                t = ct.add_table(rows=len(rows) + 1, cols=len(headers))
                no_borders(t)
                tot = sum(widths)
                for i, wd in enumerate(widths):
                    for r_ in t.rows:
                        r_.cells[i].width = Inches(wd / tot * width)
                for i, h in enumerate(headers):
                    c = t.cell(0, i)
                    shade(c, hx(NAVY))
                    nopad(c)
                    clear_cell(c)
                    p = c.add_paragraph()
                    p.paragraph_format.space_after = Pt(0)
                    if i in center:
                        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    r = p.add_run(h)
                    r.font.size = Pt(fs)
                    r.font.name = "Trebuchet MS"
                    r.font.bold = True
                    r.font.color.rgb = RGBColor(*WHITE)
                for ri, row in enumerate(rows):
                    for ci, cellv in enumerate(row):
                        c = t.cell(ri + 1, ci)
                        shade(c, hx(ROW_A if ri % 2 == 0 else ROW_B))
                        nopad(c)
                        clear_cell(c)
                        p = c.add_paragraph()
                        p.paragraph_format.space_after = Pt(0)
                        p.paragraph_format.line_spacing = 1.0
                        if ci in center:
                            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        for tx, bd in parse_runs(cellv):
                            r = p.add_run(tx)
                            r.font.size = Pt(fs)
                            r.font.name = "Courier New" if ci in mono else "Georgia"
                            r.font.bold = bool(int(bd) & BOLD) or ci == 0
                            r.font.italic = bool(int(bd) & ITAL)
                            r.font.color.rgb = RGBColor(*(NAVY if ci == 0 else TEXT))
                para(ct, "", after=4)
            elif k == "split":
                left, right, frac = b[1], b[2], b[3]
                t = ct.add_table(rows=1, cols=2)
                no_borders(t)
                lw = (width - 0.30) * frac
                rw = width - 0.30 - lw
                lc, rc = t.cell(0, 0), t.cell(0, 1)
                lc.width = Inches(lw)
                rc.width = Inches(rw)
                for cell, sub, subw in ((lc, left, lw), (rc, right, rw)):
                    nopad(cell)
                    clear_cell(cell)
                    emit(cell, sub, subw)
                para(ct, "", after=4)

    n = 1
    for sl in C.SLIDES:
        if sl.get("kind") == "title":
            para(doc, C.TITLE, size=30, color=NAVY, fam="Trebuchet MS", bold=True,
                 before=54, after=8)
            para(doc, C.SUBTITLE, size=13.5, color=MUTED, after=20)
            para(doc, "FINAL YEAR ENGINEERING PROJECT  \u00b7  DEFENSE PRESENTATION",
                 size=10, color=BLUE, fam="Trebuchet MS", bold=True, after=18)
            para(doc, "PRESENTED BY", size=9, color=BLUE, fam="Trebuchet MS",
                 bold=True, after=4)
            for nm, rl in C.TEAM:
                para(doc, f"{nm}   \u2014   {rl}", size=12, after=2)
            para(doc, "UNDER THE SUPERVISION OF", size=9, color=BLUE,
                 fam="Trebuchet MS", bold=True, before=14, after=4)
            para(doc, C.SUPERVISOR, size=12, after=6)
            para(doc, C.DEPT, size=11, color=MUTED, after=2)
            para(doc, C.COLLEGE, size=11, color=MUTED, after=2)
            para(doc, C.DATE, size=11, color=MUTED, after=2)
            continue
        n += 1
        doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
        if sl.get("kicker"):
            para(doc, f"{sl['n']}  \u00b7  {sl['kicker'].upper()}", size=9, color=BLUE,
                 fam="Trebuchet MS", bold=True, after=2)
        para(doc, sl["title"], size=19, color=NAVY, fam="Trebuchet MS", bold=True, after=3)
        if sl.get("sub"):
            para(doc, sl["sub"], size=11, color=MUTED, after=8)
        emit(doc, sl["blocks"])
    doc.save(path)
    return n


if __name__ == "__main__":
    npdf = render_pdf(OUT_PDF)
    ndoc = render_docx(OUT_DOCX)
    nppt = render_pptx(OUT_PPTX)
    print(f"PDF  {OUT_PDF}   {npdf} slides")
    print(f"DOCX {OUT_DOCX}  {ndoc} slides")
    print(f"PPTX {OUT_PPTX}  {nppt} slides")
    # report the auto-fit scale used per slide
    print("\nfit scales:")
    for i, sl in enumerate(C.SLIDES):
        if sl.get("kind") == "title":
            continue
        _, s, (bot, top) = layout_slide(sl)
        flag = "" if bot <= PH - 0.55 else "   <-- OVERFLOW"
        print(f"  {i+1:>2}  scale {s:.2f}  bottom {bot:.2f}{flag}  {sl['title'][:46]}")
